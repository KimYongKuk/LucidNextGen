"""ChromaDB 파일 업로드 & 검색"""
import os
import sys
import time
import uuid
import asyncio
import threading
from concurrent.futures import ThreadPoolExecutor
from typing import List, Dict, Optional
import chromadb

# ============================================================================
# PyTorch 2.6+/2.7+ 호환성: meta tensor 비활성화 (모델 import 전에 설정)
# ============================================================================
os.environ["TRANSFORMERS_OFFLINE"] = "0"
os.environ["HF_HUB_DISABLE_SYMLINKS_WARNING"] = "1"


def _log_timing(label: str, start_time: float, extra: str = ""):
    """타이밍 로그 출력 (stderr로 출력하여 MCP JSONRPC 스트림 보호)"""
    elapsed = time.time() - start_time
    extra_str = f" | {extra}" if extra else ""
    print(f"[TIMING] {label}: {elapsed:.3f}s{extra_str}", file=sys.stderr)
    sys.stderr.flush()
from chromadb.config import Settings
from langchain_text_splitters import RecursiveCharacterTextSplitter
import PyPDF2
from docx import Document
import openpyxl
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# PPTX 이미지 OCR 설정
PPTX_IMAGE_MIN_SIZE_BYTES = 10 * 1024  # 10KB 미만 이미지는 로고/아이콘으로 간주하여 OCR 스킵
PPTX_MAX_IMAGES_PER_FILE = 20          # 파일당 최대 OCR 대상 이미지 수 (비용 제한)
from app.services.pdf_vision_service import get_pdf_vision_service
from chromadb.utils.embedding_functions import SentenceTransformerEmbeddingFunction
from sentence_transformers import SentenceTransformer
import torch

# 스레드 풀 (ChromaDB 동기 작업을 비동기로 실행)
# CPU 코어 수에 따라 스케일링 (최대 8개)
_executor = ThreadPoolExecutor(
    max_workers=min((os.cpu_count() or 4) * 2, 8),
    thread_name_prefix="chromadb_"
)

# 모델 로드 락 (동시 로드 방지)
_model_load_lock = threading.Lock()


class SafeSentenceTransformerEmbeddingFunction:
    """PyTorch 2.7+ 호환성을 위한 안전한 SentenceTransformer 래퍼"""
    def __init__(self, model_name: str = "dragonkue/BGE-m3-ko"):
        self.model_name = model_name
        self._model = None
        self._device = None

    def name(self) -> str:
        """ChromaDB가 요구하는 name 메서드"""
        return "safe_sentence_transformer"

    def _load_model(self):
        """지연 로딩: 첫 호출 시에만 모델 로드 (동시 로드 방지 락 사용)"""
        # 이미 로드되어 있으면 바로 반환 (락 없이 빠른 경로)
        if self._model is not None:
            return self._model

        # 동시에 여러 요청이 모델을 로드하지 않도록 락 사용
        with _model_load_lock:
            # 락 획득 후 다시 확인 (다른 스레드가 이미 로드했을 수 있음)
            if self._model is not None:
                return self._model

            # GPU 사용 가능 여부 확인
            use_gpu = torch.cuda.is_available()
            if use_gpu:
                self._device = "cuda"
                print(f"[ChromaDB] CUDA available! Using GPU: {torch.cuda.get_device_name(0)}", file=sys.stderr)
            else:
                self._device = "cpu"
                print(f"[ChromaDB] CUDA not available, using CPU (this will be slow)", file=sys.stderr)

            # PyTorch 2.6+/2.7+ 호환성: GPU 직접 로드
            try:
                print(f"[ChromaDB] Loading model directly on {self._device}...", file=sys.stderr)
                self._model = SentenceTransformer(
                    self.model_name,
                    device=self._device,
                    trust_remote_code=True,
                    model_kwargs={
                        "low_cpu_mem_usage": False,  # meta tensor 사용 방지
                    }
                )
                print(f"[ChromaDB] Loaded model: {self.model_name} on {self._device}", file=sys.stderr)
            except Exception as e:
                print(f"[ChromaDB] Failed to load {self.model_name} (attempt 1): {e}", file=sys.stderr)
                # 2차 시도: torch_dtype 추가
                try:
                    self._model = SentenceTransformer(
                        self.model_name,
                        device=self._device,
                        trust_remote_code=True,
                        model_kwargs={
                            "low_cpu_mem_usage": False,
                            "torch_dtype": torch.float32,
                        }
                    )
                    print(f"[ChromaDB] Loaded model: {self.model_name} on {self._device} (attempt 2)", file=sys.stderr)
                except Exception as e2:
                    print(f"[ChromaDB] Failed to load {self.model_name} (attempt 2): {e2}", file=sys.stderr)
                    raise RuntimeError(
                        f"BGE-m3-ko 모델 로드 실패. 기존 컬렉션과 임베딩 차원 불일치를 방지하기 위해 "
                        f"폴백하지 않습니다. 원인: {e2}"
                    )
            sys.stdout.flush()
        return self._model

    def __call__(self, input: List[str]) -> List[List[float]]:
        """ChromaDB embedding function 인터페이스 구현"""
        t0 = time.time()
        model = self._load_model()
        load_time = time.time() - t0

        t1 = time.time()
        # 배치 처리 최적화: show_progress_bar=False로 오버헤드 감소
        embeddings = model.encode(
            input,
            convert_to_numpy=True,
            show_progress_bar=False,
            batch_size=32,  # 배치 크기 최적화
            normalize_embeddings=True  # 정규화 (검색 품질 향상)
        )
        encode_time = time.time() - t1

        print(f"[EMBEDDING] {len(input)} texts | device: {self._device} | model_load: {load_time:.3f}s | encode: {encode_time:.3f}s | per_text: {encode_time/len(input):.3f}s", file=sys.stderr)
        sys.stderr.flush()

        return embeddings.tolist()


class ChromaDBService:
    def __init__(self, data_path: Optional[str] = None):
        # data_path 주입을 허용해 사용자/관리자 저장소를 분리할 수 있게 함
        base_path = data_path or "./data/chromadb_user"
        abs_path = os.path.abspath(base_path)
        os.makedirs(abs_path, exist_ok=True)

        self.client = chromadb.PersistentClient(
            path=abs_path,
            settings=Settings(allow_reset=True)
        )

        # Korean-optimized embedding model (BGE-m3-ko)
        # PyTorch 2.7+ 호환성을 위한 커스텀 래퍼 사용
        self.embedding_function = SafeSentenceTransformerEmbeddingFunction(
            model_name="dragonkue/BGE-m3-ko"
        )

        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200
        )

    def get_collection(self, user_id: str, session_id: Optional[str] = None, collection: Optional[str] = None):
        """컬렉션 가져오기 (세션별, user별, 혹은 명시적 컬렉션 이름)"""
        if collection:
            name = collection
        else:
            name = f"session_{session_id}" if session_id else f"user_{user_id}"

        try:
            return self.client.get_or_create_collection(
                name,
                embedding_function=self.embedding_function
            )
        except ValueError as e:
            # 기존 컬렉션이 다른 임베딩으로 만들어져 있을 때 충돌
            # 자동 삭제 대신 에러 발생 (데이터 손실 방지)
            if "embedding function already exists" in str(e).lower():
                print(f"[ChromaDB] ERROR: Collection '{name}' has incompatible embedding model.")
                print(f"[ChromaDB] Manual migration required. Existing data will NOT be deleted automatically.")
                raise ValueError(
                    f"컬렉션 '{name}'의 임베딩 모델이 현재 설정과 호환되지 않습니다. "
                    f"데이터 보호를 위해 자동 삭제하지 않습니다. 관리자에게 문의하세요."
                )
            raise

    @staticmethod
    def _extract_table_text(table) -> str:
        """PPTX 표에서 마크다운 테이블 형식으로 텍스트 추출"""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if not rows:
            return ""
        # 첫 행 아래에 마크다운 구분선 삽입
        header = rows[0]
        separator = "| " + " | ".join("---" for _ in table.rows[0].cells) + " |"
        return header + "\n" + separator + "\n" + "\n".join(rows[1:])

    @staticmethod
    def _extract_chart_text(chart) -> str:
        """PPTX 차트에서 데이터(카테고리, 시리즈, 값) 텍스트 추출"""
        try:
            parts = []
            # 차트 제목
            if chart.has_title and chart.chart_title.has_text_frame:
                parts.append(f"[차트: {chart.chart_title.text_frame.text}]")
            else:
                parts.append("[차트]")
            # 시리즈 데이터
            for plot in chart.plots:
                categories = [str(c) for c in (plot.categories or [])]
                if categories:
                    parts.append(f"카테고리: {', '.join(categories)}")
                for series in plot.series:
                    name = series.tx.strRef.strCache[0].v if series.tx and series.tx.strRef else "데이터"
                    values = [str(v) for v in (series.values or [])]
                    parts.append(f"{name}: {', '.join(values)}")
            return "\n".join(parts)
        except Exception:
            return "[차트]"

    @staticmethod
    def _extract_shapes_text_sync(shapes, images_list: list) -> str:
        """PPTX shape들을 재귀적으로 순회하여 텍스트 추출, 이미지는 리스트에 수집

        Args:
            shapes: shape 이터러블 (slide.shapes 또는 group_shape.shapes)
            images_list: 이미지 바이트를 수집할 리스트 (mutable)

        Returns:
            추출된 텍스트
        """
        parts = []
        for shape in shapes:
            try:
                # 표
                if shape.has_table:
                    parts.append(ChromaDBService._extract_table_text(shape.table))
                    continue
                # 차트
                if hasattr(shape, "has_chart") and shape.has_chart:
                    parts.append(ChromaDBService._extract_chart_text(shape.chart))
                    continue
                # 그룹 shape → 재귀 탐색
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    parts.append(ChromaDBService._extract_shapes_text_sync(shape.shapes, images_list))
                    continue
                # 이미지
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        blob = shape.image.blob
                        if len(blob) >= PPTX_IMAGE_MIN_SIZE_BYTES and len(images_list) < PPTX_MAX_IMAGES_PER_FILE:
                            images_list.append(blob)
                    except Exception:
                        pass  # 외부 링크 이미지 등 blob 접근 불가
                    continue
                # 일반 텍스트
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
            except Exception as e:
                print(f"[PPTX] Shape 처리 오류: {e}", file=sys.stderr)
                continue
        return "\n".join(parts)

    def _extract_excel_sheets(self, file_path: str) -> List[Dict]:
        """엑셀 파일에서 시트별 구조화된 데이터 추출

        Returns:
            [{"sheet_name": str, "header": str, "rows": List[str]}]
        """
        wb = openpyxl.load_workbook(file_path, data_only=True)
        sheets_data = []

        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell)
                if row_text.strip():
                    rows.append(row_text)

            if not rows:
                continue

            sheets_data.append({
                "sheet_name": sheet.title,
                "header": rows[0],
                "rows": rows,
            })

        return sheets_data

    def _chunk_excel_sheets(self, sheets_data: List[Dict], chunk_size: int = 1000) -> List[Dict]:
        """시트별 행 기반 청킹 (시트명 + 헤더 프리픽스 포함)

        Returns:
            [{"text": str, "sheet_name": str}]
        """
        chunks = []

        for sheet in sheets_data:
            sheet_name = sheet["sheet_name"]
            header = sheet["header"]
            rows = sheet["rows"]

            prefix = f"[시트: {sheet_name}]\n{header}\n"
            prefix_len = len(prefix)

            # 데이터 행 (헤더 제외)
            data_rows = rows[1:] if len(rows) > 1 else []

            if not data_rows:
                # 헤더만 있는 시트
                chunks.append({"text": prefix.rstrip("\n"), "sheet_name": sheet_name})
                

            current_chunk_rows = []
            current_len = prefix_len

            for row in data_rows:
                row_len = len(row) + 1  # +1 for newline

                if current_len + row_len > chunk_size and current_chunk_rows:
                    chunk_text = prefix + "\n".join(current_chunk_rows)
                    chunks.append({"text": chunk_text, "sheet_name": sheet_name})
                    current_chunk_rows = []
                    current_len = prefix_len

                current_chunk_rows.append(row)
                current_len += row_len

            # 마지막 청크
            if current_chunk_rows:
                chunk_text = prefix + "\n".join(current_chunk_rows)
                chunks.append({"text": chunk_text, "sheet_name": sheet_name})

        return chunks

    async def extract_text(self, file_path: str, filename: str) -> str:
        """파일에서 텍스트 추출 (PDF는 하이브리드 처리)"""
        ext = filename.lower().split('.')[-1]

        if ext == 'pdf':
            # 하이브리드 PDF 처리 (텍스트 + 이미지)
            pdf_vision = get_pdf_vision_service()
            page_results = await pdf_vision.process_pdf(file_path)
            return pdf_vision.combine_page_contents(page_results)

        elif ext in ['docx', 'doc']:
            doc = Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs)

        elif ext in ['xlsx', 'xls']:
            wb = openpyxl.load_workbook(file_path)
            text = ""
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    text += " | ".join(str(cell) for cell in row if cell) + "\n"
            return text

        elif ext in ['pptx', 'ppt']:
            prs = Presentation(file_path)
            all_text_parts = []
            all_images = []  # (slide_idx, image_bytes) 튜플 리스트

            for slide_idx, slide in enumerate(prs.slides):
                all_text_parts.append(f"\n--- 슬라이드 {slide_idx + 1} ---\n")
                slide_images = []
                text = self._extract_shapes_text_sync(slide.shapes, slide_images)
                all_text_parts.append(text)
                all_images.extend([(slide_idx, img) for img in slide_images])

            # 이미지 Vision API 일괄 OCR 처리
            if all_images:
                t0 = time.time()
                pdf_vision = get_pdf_vision_service()
                semaphore = asyncio.Semaphore(pdf_vision.vision_concurrency)

                async def ocr_image(slide_idx, img_bytes):
                    async with semaphore:
                        return (slide_idx, await pdf_vision.extract_text_from_image(img_bytes))

                print(f"[PPTX] Vision API OCR: {len(all_images)}개 이미지 처리 시작", file=sys.stderr)
                results = await asyncio.gather(
                    *[ocr_image(si, ib) for si, ib in all_images],
                    return_exceptions=True
                )

                # 슬라이드별 OCR 텍스트 병합
                ocr_by_slide: Dict[int, List[str]] = {}
                for r in results:
                    if isinstance(r, tuple):
                        si, ocr_text = r
                        if ocr_text and ocr_text.strip():
                            ocr_by_slide.setdefault(si, []).append(ocr_text)

                for si in sorted(ocr_by_slide.keys()):
                    all_text_parts.append(f"\n[슬라이드 {si + 1} 이미지 내용]\n")
                    all_text_parts.append("\n".join(ocr_by_slide[si]))

                _log_timing("PPTX Vision OCR", t0, f"{len(all_images)} images")

            return "\n".join(all_text_parts)

        elif ext == 'txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()

        elif ext in ['html', 'htm']:
            from bs4 import BeautifulSoup
            import re, html as html_mod
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            soup = BeautifulSoup(html_content, 'html.parser')

            # Extract structured data from script tags (e.g. email backup metadata)
            script_data_lines = []
            folder_name = ""
            for script_tag in soup.find_all('script'):
                script_text = script_tag.string or ""
                # Extract folder name
                folder_match = re.search(r'folder\s*=\s*"([^"]*)"', script_text)
                if folder_match:
                    folder_name = folder_match.group(1)
                # Match JS object literals like li[0] = {"subject":"...", "from":"...", ...}
                entries = re.findall(r'\{[^{}]*"subject"\s*:\s*"[^"]*"[^{}]*\}', script_text)
                for entry in entries:
                    pairs = re.findall(r'"(\w+)"\s*:\s*"([^"]*)"', entry)
                    if pairs:
                        parts = []
                        for key, val in pairs:
                            if key in ('attach', 'html', 'size'):
                                continue
                            decoded = html_mod.unescape(val).strip()
                            if decoded:
                                parts.append(f"{key}: {decoded}")
                        if parts:
                            script_data_lines.append(" | ".join(parts))

            # Remove script and style tags for body text
            for tag in soup(['script', 'style']):
                tag.decompose()
            body_text = soup.get_text(separator='\n', strip=True)

            if script_data_lines:
                # Return with special separator for row-based chunking
                prefix = f"[메일함: {folder_name}] " if folder_name else ""
                header = f"{prefix}메일 목록 ({len(script_data_lines)}건)"
                return "<<EMAIL_LIST>>\n" + header + "\n" + "\n".join(script_data_lines) + "\n\n" + body_text
            return body_text

        elif ext == 'csv':
            import pandas as pd
            df = pd.read_csv(file_path, encoding='utf-8')
            return df.to_string()

        raise ValueError(f"지원하지 않는 파일 형식: {ext}")

    def _sync_add_to_collection(
        self,
        collection_obj,
        ids: List[str],
        chunks: List[str],
        metadatas: List[Dict],
    ):
        """ChromaDB 동기 쓰기 작업 (스레드 풀에서 실행)"""
        collection_obj.add(
            ids=ids,
            documents=chunks,
            metadatas=metadatas
        )

    async def upload_file(
        self,
        file_content: bytes,
        filename: str,
        user_id: str,
        session_id: Optional[str] = None,
        replace_existing: bool = False,
        collection: Optional[str] = None,
        chunk_size: Optional[int] = None,
        chunk_overlap: Optional[int] = None,
        file_id: Optional[str] = None,
    ) -> Dict:
        """
        파일 업로드 및 벡터화 (세션별 또는 user별)

        Args:
            file_content: 파일 내용
            filename: 파일명
            user_id: 사용자 ID
            session_id: 세션 ID (임시 파일용)
            replace_existing: True면 기존 파일 삭제 후 업로드 (기본값: False, 여러 파일 업로드 지원)
        """
        upload_start = time.time()
        print(f"\n{'='*60}", file=sys.stderr)
        print(f"[UPLOAD START] {filename} ({len(file_content)} bytes)", file=sys.stderr)
        print(f"{'='*60}", file=sys.stderr)
        sys.stderr.flush()

        # 세션별 업로드인 경우, replace_existing이 명시적으로 True일 때만 기존 컬렉션 삭제
        # (명시적 컬렉션 지정 시에는 건드리지 않음)
        # 기본값을 False로 변경하여 여러 파일 동시 업로드를 지원
        if session_id and replace_existing and not collection:
            # 기존 세션 컬렉션 삭제 (새 파일로 교체)
            try:
                await asyncio.get_event_loop().run_in_executor(
                    _executor,
                    self.client.delete_collection,
                    f"session_{session_id}"
                )
            except Exception:
                # 컬렉션이 없으면 무시
                pass

        # 임시 파일 저장
        t0 = time.time()
        # 임시 파일명에 원본 파일명을 포함하면 한글/괄호/공백 등으로 경로 문제 발생
        # 확장자만 유지하여 안전한 경로 생성
        ext_for_temp = filename.rsplit('.', 1)[-1] if '.' in filename else 'tmp'
        temp_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), f"temp_{uuid.uuid4()}.{ext_for_temp}")

        # aiofiles가 설치되어 있다고 가정하거나, run_in_executor 사용
        # 여기서는 표준 라이브러리만 사용하기 위해 run_in_executor 사용
        def write_temp_file():
            with open(temp_path, 'wb') as f:
                f.write(file_content)

        await asyncio.get_event_loop().run_in_executor(
            None, # default executor
            write_temp_file
        )
        _log_timing("Write temp file", t0)

        try:
            ext = filename.lower().rsplit('.', 1)[-1]

            if ext in ['xlsx', 'xls']:
                # ── 엑셀: 시트별 행 기반 청킹 (헤더 프리픽스) ──
                t0 = time.time()
                sheets_data = self._extract_excel_sheets(temp_path)
                total_rows = sum(len(s["rows"]) for s in sheets_data)
                _log_timing("Excel extraction", t0, f"{len(sheets_data)} sheets, {total_rows} rows")

                t0 = time.time()
                target_chunk_size = chunk_size or self.splitter._chunk_size
                excel_chunks = self._chunk_excel_sheets(sheets_data, target_chunk_size)
                chunks = [c["text"] for c in excel_chunks]
                sheet_names = [c["sheet_name"] for c in excel_chunks]
                _log_timing("Excel chunking (row-based)", t0, f"{len(chunks)} chunks")
            else:
                # ── 기존 방식: 텍스트 추출 → RecursiveCharacterTextSplitter ──
                t0 = time.time()
                text = await self.extract_text(temp_path, filename)
                _log_timing("Text extraction (total)", t0, f"{len(text)} chars")

                t0 = time.time()
                # Email HTML: row-based chunking to prevent splitting email entries
                if text.startswith("<<EMAIL_LIST>>\n"):
                    text = text[len("<<EMAIL_LIST>>\n"):]
                    lines = text.split("\n")
                    header_line = lines[0] if lines else ""  # e.g. "[메일함: 보낸메일함] 메일 목록 (17건)"
                    target = chunk_size or self.splitter._chunk_size
                    chunks = []
                    current_chunk_lines = []
                    current_size = len(header_line) + 1
                    for line in lines[1:]:
                        line_size = len(line) + 1
                        if current_chunk_lines and current_size + line_size > target:
                            chunks.append(header_line + "\n" + "\n".join(current_chunk_lines))
                            current_chunk_lines = []
                            current_size = len(header_line) + 1
                        current_chunk_lines.append(line)
                        current_size += line_size
                    if current_chunk_lines:
                        chunks.append(header_line + "\n" + "\n".join(current_chunk_lines))
                else:
                    splitter = self.splitter
                    if chunk_size or chunk_overlap:
                        splitter = RecursiveCharacterTextSplitter(
                            chunk_size=chunk_size or self.splitter._chunk_size,
                            chunk_overlap=chunk_overlap or self.splitter._chunk_overlap,
                        )
                    chunks = splitter.split_text(text)
                _log_timing("Chunking", t0, f"{len(chunks)} chunks")
                sheet_names = None

            # ChromaDB에 저장 (세션별 또는 user별) - 스레드 풀에서 실행
            t0 = time.time()
            collection_obj = self.get_collection(user_id, session_id, collection=collection)
            _log_timing("Get collection", t0)

            if not file_id:
                file_id = str(uuid.uuid4())

            # 업로드 시간 기록 (ISO format)
            from datetime import datetime, timezone
            uploaded_at = datetime.now(timezone.utc).isoformat()

            ids = [f"{file_id}_{i}" for i in range(len(chunks))]
            metadatas = []
            for i in range(len(chunks)):
                meta = {"filename": filename, "file_id": file_id, "uploaded_at": uploaded_at}
                if sheet_names and sheet_names[i]:
                    meta["sheet_name"] = sheet_names[i]
                metadatas.append(meta)

            # ✅ ChromaDB 쓰기를 별도 스레드에서 실행 (다른 요청 차단 안함)
            t0 = time.time()
            await asyncio.get_event_loop().run_in_executor(
                _executor,
                self._sync_add_to_collection,
                collection_obj,
                ids,
                chunks,
                metadatas
            )
            _log_timing("ChromaDB add (embedding + write)", t0, f"{len(chunks)} chunks")

            print(f"\n{'='*60}", file=sys.stderr)
            _log_timing("UPLOAD TOTAL", upload_start)
            print(f"{'='*60}\n", file=sys.stderr)
            sys.stderr.flush()

            return {
                "success": True,
                "file_id": file_id,
                "filename": filename,
                "chunks": len(chunks)
            }

        finally:
            # 임시 파일 삭제
            if os.path.exists(temp_path):
                os.remove(temp_path)

    def _sync_search_collection(
        self,
        collection_obj,
        query: str,
        limit: int,
        min_relevance: float = 0.0,
    ) -> List[Dict]:
        """ChromaDB 동기 검색 작업 (스레드 풀에서 실행)

        Args:
            collection_obj: ChromaDB 컬렉션
            query: 검색 쿼리
            limit: 최대 결과 수
            min_relevance: 최소 유사도 임계값 (0.0~1.0, 기본값 0.0 = 필터링 없음)
                          ChromaDB 코사인 거리는 0~2 범위이므로 내부적으로 변환
        """
        if collection_obj.count() == 0:
            return []

        results = collection_obj.query(
            query_texts=[query],
            n_results=limit,
            include=["documents", "metadatas", "distances"]  # 거리 정보 포함
        )

        docs = []
        if results['documents'] and results['documents'][0]:
            for i, doc in enumerate(results['documents'][0]):
                # 코사인 거리 → 유사도 변환 (distance: 0~2, similarity: 1~-1)
                # similarity = 1 - (distance / 2) 로 0~1 범위로 정규화
                distance = results['distances'][0][i] if results.get('distances') else 0
                similarity = 1 - (distance / 2)  # 0~1 범위

                # 유사도 임계값 필터링
                if min_relevance > 0 and similarity < min_relevance:
                    print(f"[FILTER] Skipped (similarity={similarity:.3f} < {min_relevance}): {doc[:50]}...", file=sys.stderr)
                    continue

                docs.append({
                    "text": doc,
                    "metadata": results['metadatas'][0][i] if results['metadatas'] else {},
                    "similarity": round(similarity, 3)  # 유사도 점수 포함
                })

        return docs

    async def search(
        self,
        query: str,
        user_id: str,
        session_id: Optional[str] = None,
        limit: int = 5,
        collection: Optional[str] = None,
        min_relevance: float = 0.0,
    ) -> List[Dict]:
        """파일 검색 (세션별 또는 user별) - 비동기 실행

        Args:
            query: 검색 쿼리
            user_id: 사용자 ID
            session_id: 세션 ID
            limit: 최대 결과 수
            collection: 명시적 컬렉션 이름
            min_relevance: 최소 유사도 임계값 (0.0~1.0)
        """
        collection_obj = self.get_collection(user_id, session_id, collection=collection)

        # ✅ ChromaDB 검색을 별도 스레드에서 실행 (다른 요청 차단 안함)
        # functools.partial 사용하여 키워드 인자 전달
        import functools
        search_func = functools.partial(
            self._sync_search_collection,
            collection_obj,
            query,
            limit,
            min_relevance
        )
        docs = await asyncio.get_event_loop().run_in_executor(
            _executor,
            search_func
        )

        return docs

    def has_session_files(self, session_id: str) -> bool:
        """세션에 업로드된 파일이 있는지 확인"""
        try:
            collection = self.get_collection("anonymous", session_id)
            return collection.count() > 0
        except Exception:
            return False

    async def delete_session_files(self, session_id: str) -> Dict:
        """세션 파일 전체 삭제"""
        try:
            self.client.delete_collection(f"session_{session_id}")
            return {
                "success": True,
                "message": f"Session {session_id} files deleted"
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e)
            }


# 싱글톤 캐시 (사용자/관리자 분리)
_chromadb_service_user: Optional[ChromaDBService] = None
_chromadb_service_admin: Optional[ChromaDBService] = None


def get_chromadb_service() -> ChromaDBService:
    """기존 호환용: 기본적으로 사용자 저장소(chromadb_user)를 사용."""
    return get_user_chromadb_service()


def get_user_chromadb_service() -> ChromaDBService:
    global _chromadb_service_user
    if _chromadb_service_user is None:
        _chromadb_service_user = ChromaDBService(data_path="./data/chromadb_user")
    return _chromadb_service_user


def get_admin_chromadb_service() -> ChromaDBService:
    global _chromadb_service_admin
    if _chromadb_service_admin is None:
        _chromadb_service_admin = ChromaDBService(data_path="./data/chromadb_admin")
    return _chromadb_service_admin
