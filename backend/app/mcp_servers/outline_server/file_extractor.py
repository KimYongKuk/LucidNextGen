"""파일 → 마크다운 + 이미지 추출 모듈

PDF / PPTX / DOCX 파일에서 텍스트를 마크다운으로 변환하고,
임베디드 이미지를 바이너리로 추출하여 스테이징 디렉토리에 저장합니다.

추출 전략 (C안):
  - 텍스트 → 마크다운 (헤딩, 목록, 본문)
  - 임베디드 이미지 → 바이너리 추출 후 스테이징
  - 표 → 마크다운 테이블
  - 벡터 다이어그램 → 미지원 (추후 보완)
"""

import os
import sys
import uuid
import json
import base64
import mimetypes
from pathlib import Path
from typing import Dict, List, Optional, Any

# ── 상수 ──────────────────────────────────────────────
IMAGE_MIN_SIZE_BYTES = 5 * 1024       # 5KB 미만 이미지 필터 (아이콘/로고)
IMAGE_MAX_COUNT = 15                   # 문서당 최대 이미지 수
MARKDOWN_MAX_LENGTH = 100_000          # 마크다운 최대 길이 (Outline 문서 한도)
IMAGE_DESCRIPTION_MAX_TOKENS = 300     # 이미지 설명 최대 토큰

# 스테이징 루트 (server.py에서 설정 가능)
STAGING_ROOT = Path(os.environ.get(
    "OUTLINE_STAGING_DIR",
    str(Path(__file__).parent.parent.parent.parent / "data" / "outline_staging"),
))


# ── 공통 유틸 ─────────────────────────────────────────

def _guess_image_ext(image_bytes: bytes, fallback: str = "png") -> str:
    """이미지 바이너리에서 확장자 추측"""
    if image_bytes[:8] == b'\x89PNG\r\n\x1a\n':
        return "png"
    if image_bytes[:2] == b'\xff\xd8':
        return "jpeg"
    if image_bytes[:4] == b'GIF8':
        return "gif"
    if image_bytes[:4] == b'RIFF' and image_bytes[8:12] == b'WEBP':
        return "webp"
    return fallback


def _save_image(
    image_bytes: bytes,
    staging_dir: Path,
    index: int,
    ext: str = "",
) -> Optional[Dict[str, str]]:
    """이미지를 스테이징 디렉토리에 저장하고 메타 반환"""
    if len(image_bytes) < IMAGE_MIN_SIZE_BYTES:
        return None

    if not ext:
        ext = _guess_image_ext(image_bytes)

    filename = f"image_{index}.{ext}"
    path = staging_dir / filename
    path.write_bytes(image_bytes)

    content_type = mimetypes.types_map.get(f".{ext}", "image/png")

    return {
        "placeholder": f"{{{{IMAGE_{index}}}}}",
        "filename": filename,
        "path": str(path),
        "content_type": content_type,
        "size": len(image_bytes),
    }


# ── Vision API (이미지 설명 생성) ──────────────────────

def _describe_image_via_vision(image_bytes: bytes, content_type: str = "image/png") -> str:
    """Bedrock Claude Vision API로 이미지 내용 설명 생성 (동기 호출)

    정제 모드에서만 사용. 이미지 안의 텍스트, 다이어그램, 표 등을 설명합니다.
    MCP 서버 서브프로세스에서 직접 Bedrock 호출.
    """
    try:
        import boto3
        from botocore.config import Config as BotoConfig
    except ImportError:
        return ""

    b64 = base64.b64encode(image_bytes).decode("utf-8")

    # media_type 정규화
    media_map = {"image/jpg": "image/jpeg"}
    media_type = media_map.get(content_type, content_type)
    if media_type not in ("image/png", "image/jpeg", "image/gif", "image/webp"):
        media_type = "image/png"

    request_body = {
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": IMAGE_DESCRIPTION_MAX_TOKENS,
        "temperature": 0,
        "messages": [{
            "role": "user",
            "content": [
                {
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": media_type,
                        "data": b64,
                    },
                },
                {
                    "type": "text",
                    "text": (
                        "이 이미지의 내용을 한국어로 간결하게 설명하세요. "
                        "다이어그램이면 흐름/구성요소를, 표면 주요 항목을, "
                        "텍스트가 포함되어 있으면 핵심 내용을 포함하세요. "
                        "2~3문장 이내로."
                    ),
                },
            ],
        }],
    }

    try:
        region = os.environ.get("AWS_REGION", "us-west-2")
        model_id = os.environ.get(
            "BEDROCK_HAIKU_MODEL_ID",
            "us.anthropic.claude-haiku-4-5-20251001-v1:0",
        )
        config = BotoConfig(read_timeout=30, connect_timeout=10)
        client = boto3.client("bedrock-runtime", region_name=region, config=config)

        resp = client.invoke_model(
            modelId=model_id,
            body=json.dumps(request_body),
            contentType="application/json",
            accept="application/json",
        )
        body = json.loads(resp["body"].read())
        if body.get("content"):
            return body["content"][0].get("text", "")
    except Exception as e:
        print(f"[FileExtractor] Vision API 오류: {e}", file=sys.stderr)
    return ""


def describe_images(images: List[Dict]) -> List[Dict]:
    """이미지 목록에 description 필드를 추가합니다 (정제 모드용).

    각 이미지의 staging path에서 바이너리를 읽어 Vision API로 설명을 생성합니다.

    Args:
        images: extract_file 결과의 images 리스트

    Returns:
        description 필드가 추가된 images 리스트
    """
    for img in images:
        img_path = Path(img.get("path", ""))
        if not img_path.exists():
            img["description"] = ""
            continue

        try:
            img_bytes = img_path.read_bytes()
            content_type = img.get("content_type", "image/png")
            description = _describe_image_via_vision(img_bytes, content_type)
            img["description"] = description
        except Exception:
            img["description"] = ""

    return images


# ── PDF 추출 ─────────────────────────────────────────

def _extract_pdf(file_path: str, staging_dir: Path) -> Dict[str, Any]:
    """PDF에서 텍스트(마크다운) + 임베디드 이미지 추출

    PyMuPDF의 get_text("dict")로 블록 순서 보존,
    폰트 크기 기반 헤딩 감지, extract_image(xref)로 이미지 바이너리 추출.
    """
    import fitz  # PyMuPDF

    doc = fitz.open(file_path)
    md_parts: List[str] = []
    images: List[Dict] = []
    image_idx = 0

    # 전체 문서의 폰트 크기 분포 분석 (헤딩 감지용)
    all_font_sizes: List[float] = []
    for page in doc:
        blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]
        for block in blocks:
            if block["type"] != 0:  # 텍스트 블록만
                continue
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    text = span.get("text", "").strip()
                    if text:
                        all_font_sizes.append(span.get("size", 12))

    # 폰트 크기 기반 헤딩 임계값 계산
    if all_font_sizes:
        avg_size = sum(all_font_sizes) / len(all_font_sizes)
        h1_threshold = avg_size * 1.6    # 60% 이상 크면 H1
        h2_threshold = avg_size * 1.3    # 30% 이상 크면 H2
        h3_threshold = avg_size * 1.15   # 15% 이상 크면 H3
    else:
        avg_size = 12
        h1_threshold = 20
        h2_threshold = 16
        h3_threshold = 14

    # 이미 추출한 xref 추적 (중복 방지)
    extracted_xrefs: set = set()

    for page_num, page in enumerate(doc):
        # ── 이미지 추출 ──
        if image_idx < IMAGE_MAX_COUNT:
            page_images = page.get_images(full=True)
            for img_info in page_images:
                if image_idx >= IMAGE_MAX_COUNT:
                    break
                xref = img_info[0]
                if xref in extracted_xrefs:
                    continue
                extracted_xrefs.add(xref)

                try:
                    img_data = doc.extract_image(xref)
                    if not img_data or not img_data.get("image"):
                        continue
                    img_bytes = img_data["image"]
                    ext = img_data.get("ext", "png")
                    meta = _save_image(img_bytes, staging_dir, image_idx, ext)
                    if meta:
                        images.append(meta)
                        image_idx += 1
                except Exception:
                    continue

        # ── 텍스트 추출 ──
        blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]
        page_text_parts: List[str] = []

        for block in blocks:
            if block["type"] != 0:  # 이미지 블록은 skip (이미 추출)
                continue

            block_text = ""
            block_max_size = 0
            is_bold = False

            for line in block.get("lines", []):
                line_text = ""
                for span in line.get("spans", []):
                    text = span.get("text", "")
                    size = span.get("size", avg_size)
                    flags = span.get("flags", 0)

                    if size > block_max_size:
                        block_max_size = size
                    if flags & 2 ** 4:  # bold flag
                        is_bold = True

                    line_text += text

                line_text = line_text.rstrip()
                if line_text:
                    block_text += line_text + "\n"

            block_text = block_text.strip()
            if not block_text:
                continue

            # 헤딩 감지 (폰트 크기 + 짧은 텍스트)
            is_short = len(block_text) < 200
            if is_short and block_max_size >= h1_threshold:
                page_text_parts.append(f"# {block_text}")
            elif is_short and block_max_size >= h2_threshold:
                page_text_parts.append(f"## {block_text}")
            elif is_short and (block_max_size >= h3_threshold or is_bold):
                page_text_parts.append(f"### {block_text}")
            else:
                page_text_parts.append(block_text)

        if page_text_parts:
            md_parts.append("\n\n".join(page_text_parts))

        # 이미지 플레이스홀더 삽입 (해당 페이지에서 추출된 이미지)
        # 페이지 끝에 배치 (정확한 위치 추적은 복잡하므로 페이지 단위)
        page_image_start = sum(1 for img in images if img in images[:image_idx])

    # 이미지 플레이스홀더를 마크다운에 삽입
    # 각 이미지는 추출 순서대로 문서 끝에 배치 (간단한 전략)
    markdown = "\n\n".join(md_parts)

    if images:
        markdown += "\n\n---\n\n"
        for img in images:
            markdown += f"\n\n![{img['filename']}]({img['placeholder']})\n"

    doc.close()
    return {"markdown": markdown, "images": images}


# ── PPTX 추출 ────────────────────────────────────────

def _extract_pptx(file_path: str, staging_dir: Path) -> Dict[str, Any]:
    """PPTX에서 슬라이드별 텍스트(마크다운) + 이미지 추출

    shape 재귀 순회로 텍스트/표/차트/이미지 추출.
    chromadb_service.py의 패턴을 참고하되 마크다운 구조 보존.
    """
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    md_parts: List[str] = []
    images: List[Dict] = []
    image_idx = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts: List[str] = []

        # 슬라이드 제목 감지
        title_text = ""
        if slide.shapes.title and slide.shapes.title.text.strip():
            title_text = slide.shapes.title.text.strip()
            slide_parts.append(f"## {title_text}")
        else:
            slide_parts.append(f"## 슬라이드 {slide_num}")

        # shape 순회
        for shape in slide.shapes:
            try:
                # 제목은 이미 처리
                if shape == slide.shapes.title:
                    continue

                result = _extract_pptx_shape(
                    shape, staging_dir, images, image_idx, MSO_SHAPE_TYPE
                )
                if result["text"]:
                    slide_parts.append(result["text"])
                image_idx = result["image_idx"]

            except Exception:
                continue

        md_parts.append("\n\n".join(slide_parts))

    markdown = "\n\n---\n\n".join(md_parts)
    return {"markdown": markdown, "images": images}


def _extract_pptx_shape(
    shape, staging_dir: Path, images: List[Dict],
    image_idx: int, MSO_SHAPE_TYPE,
) -> Dict[str, Any]:
    """단일 PPTX shape에서 텍스트/이미지 추출 (재귀)"""
    text_parts: List[str] = []

    # 표
    if shape.has_table:
        text_parts.append(_extract_pptx_table(shape.table))
        return {"text": "\n".join(text_parts), "image_idx": image_idx}

    # 차트
    if hasattr(shape, "has_chart") and shape.has_chart:
        text_parts.append(_extract_pptx_chart(shape.chart))
        return {"text": "\n".join(text_parts), "image_idx": image_idx}

    # 그룹 shape → 재귀
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            result = _extract_pptx_shape(
                child, staging_dir, images, image_idx, MSO_SHAPE_TYPE
            )
            if result["text"]:
                text_parts.append(result["text"])
            image_idx = result["image_idx"]
        return {"text": "\n".join(text_parts), "image_idx": image_idx}

    # 이미지
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        if image_idx < IMAGE_MAX_COUNT:
            try:
                blob = shape.image.blob
                meta = _save_image(blob, staging_dir, image_idx)
                if meta:
                    images.append(meta)
                    text_parts.append(f"\n![{meta['filename']}]({meta['placeholder']})\n")
                    image_idx += 1
            except Exception:
                pass  # 외부 링크 이미지 등
        return {"text": "\n".join(text_parts), "image_idx": image_idx}

    # 일반 텍스트
    if hasattr(shape, "text") and shape.text.strip():
        text_parts.append(shape.text.strip())

    return {"text": "\n".join(text_parts), "image_idx": image_idx}


def _extract_pptx_table(table) -> str:
    """PPTX 표 → 마크다운 테이블"""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
    if not rows:
        return ""
    header = rows[0]
    separator = "| " + " | ".join("---" for _ in table.rows[0].cells) + " |"
    return header + "\n" + separator + "\n" + "\n".join(rows[1:])


def _extract_pptx_chart(chart) -> str:
    """PPTX 차트 → 텍스트 표현"""
    try:
        parts = []
        if chart.has_title and chart.chart_title.has_text_frame:
            parts.append(f"**[차트: {chart.chart_title.text_frame.text}]**")
        else:
            parts.append("**[차트]**")
        for plot in chart.plots:
            categories = [str(c) for c in (plot.categories or [])]
            if categories:
                parts.append(f"카테고리: {', '.join(categories)}")
            for series in plot.series:
                name = "데이터"
                try:
                    name = series.tx.strRef.strCache[0].v
                except Exception:
                    pass
                values = [str(v) for v in (series.values or [])]
                parts.append(f"{name}: {', '.join(values)}")
        return "\n".join(parts)
    except Exception:
        return "**[차트]**"


# ── DOCX 추출 ────────────────────────────────────────

def _extract_docx(file_path: str, staging_dir: Path) -> Dict[str, Any]:
    """DOCX에서 문단/표/이미지 추출 → 마크다운

    paragraph.style.name 기반 헤딩 감지,
    inline_shape → relationship → blob으로 이미지 추출,
    doc.tables → 마크다운 표 변환.
    """
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(file_path)
    md_parts: List[str] = []
    images: List[Dict] = []
    image_idx = 0

    # 표 위치 추적 (중복 렌더링 방지)
    table_elements = set()
    for table in doc.tables:
        table_elements.add(table._tbl)

    # body의 모든 요소를 순서대로 순회
    body = doc.element.body
    for child in body:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

        # ── 문단 (paragraph) ──
        if tag == "p":
            paragraph = None
            for p in doc.paragraphs:
                if p._element is child:
                    paragraph = p
                    break
            if paragraph is None:
                continue

            text = paragraph.text.strip()
            style_name = paragraph.style.name if paragraph.style else ""

            # 헤딩 감지
            if style_name.startswith("Heading 1") or style_name == "제목 1":
                md_parts.append(f"# {text}" if text else "")
            elif style_name.startswith("Heading 2") or style_name == "제목 2":
                md_parts.append(f"## {text}" if text else "")
            elif style_name.startswith("Heading 3") or style_name == "제목 3":
                md_parts.append(f"### {text}" if text else "")
            elif style_name.startswith("Heading") or "제목" in style_name:
                md_parts.append(f"#### {text}" if text else "")
            elif style_name.startswith("List"):
                md_parts.append(f"- {text}" if text else "")
            elif text:
                md_parts.append(text)

            # 인라인 이미지 추출
            if image_idx < IMAGE_MAX_COUNT:
                for run in paragraph.runs:
                    drawing_elements = run._element.findall(qn("w:drawing"))
                    for drawing in drawing_elements:
                        blip_elements = drawing.findall(".//" + qn("a:blip"))
                        for blip in blip_elements:
                            if image_idx >= IMAGE_MAX_COUNT:
                                break
                            r_embed = blip.get(qn("r:embed"))
                            if not r_embed:
                                continue
                            try:
                                rel = doc.part.rels[r_embed]
                                img_bytes = rel.target_part.blob
                                content_type = rel.target_part.content_type or ""
                                ext = _ext_from_content_type(content_type)
                                meta = _save_image(img_bytes, staging_dir, image_idx, ext)
                                if meta:
                                    images.append(meta)
                                    md_parts.append(f"\n![{meta['filename']}]({meta['placeholder']})\n")
                                    image_idx += 1
                            except Exception:
                                continue

        # ── 표 (table) ──
        elif tag == "tbl":
            for table in doc.tables:
                if table._tbl is child:
                    md_parts.append(_extract_docx_table(table))
                    break

    markdown = "\n\n".join(part for part in md_parts if part)
    return {"markdown": markdown, "images": images}


def _ext_from_content_type(content_type: str) -> str:
    """Content-Type에서 확장자 추출"""
    mapping = {
        "image/png": "png",
        "image/jpeg": "jpeg",
        "image/gif": "gif",
        "image/webp": "webp",
        "image/bmp": "bmp",
        "image/tiff": "tiff",
        "image/svg+xml": "svg",
    }
    return mapping.get(content_type.lower(), "png")


def _extract_docx_table(table) -> str:
    """DOCX 표 → 마크다운 테이블"""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
    if not rows:
        return ""
    col_count = len(table.rows[0].cells)
    header = rows[0]
    separator = "| " + " | ".join("---" for _ in range(col_count)) + " |"
    return header + "\n" + separator + "\n" + "\n".join(rows[1:])


# ── 공개 API ─────────────────────────────────────────

SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx"}


def extract_file(file_path: str, staging_id: str = "") -> Dict[str, Any]:
    """파일에서 마크다운 + 이미지를 추출합니다.

    Args:
        file_path: 파일 경로
        staging_id: 스테이징 디렉토리 ID (없으면 UUID 자동 생성)

    Returns:
        {
            "title": str,           # 파일명 (확장자 제외)
            "markdown": str,        # 마크다운 본문 (이미지는 플레이스홀더)
            "images": [             # 추출된 이미지 목록
                {
                    "placeholder": "{{IMAGE_0}}",
                    "filename": "image_0.png",
                    "path": "/absolute/path/to/staged/image_0.png",
                    "content_type": "image/png",
                    "size": 12345,
                }
            ],
            "staging_dir": str,     # 스테이징 디렉토리 경로
        }
    """
    path = Path(file_path)
    if not path.exists():
        return {"error": f"파일을 찾을 수 없습니다: {file_path}"}

    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        return {"error": f"지원하지 않는 형식입니다: {ext} (지원: PDF, PPTX, DOCX)"}

    # 스테이징 디렉토리 생성
    if not staging_id:
        staging_id = uuid.uuid4().hex[:12]
    staging_dir = STAGING_ROOT / staging_id
    staging_dir.mkdir(parents=True, exist_ok=True)

    # 포맷별 추출
    try:
        if ext == ".pdf":
            result = _extract_pdf(file_path, staging_dir)
        elif ext == ".pptx":
            result = _extract_pptx(file_path, staging_dir)
        elif ext == ".docx":
            result = _extract_docx(file_path, staging_dir)
        else:
            return {"error": f"지원하지 않는 형식: {ext}"}
    except Exception as e:
        return {"error": f"파일 추출 실패: {str(e)[:500]}"}

    # 마크다운 길이 제한
    markdown = result.get("markdown", "")
    if len(markdown) > MARKDOWN_MAX_LENGTH:
        markdown = markdown[:MARKDOWN_MAX_LENGTH].rstrip()
        markdown += f"\n\n> [본문이 {MARKDOWN_MAX_LENGTH:,}자로 잘렸습니다]"

    title = path.stem  # 확장자 제외 파일명

    return {
        "title": title,
        "markdown": markdown,
        "images": result.get("images", []),
        "staging_dir": str(staging_dir),
    }


def cleanup_old_staging(max_age_hours: int = 1) -> int:
    """오래된 스테이징 디렉토리 정리

    Returns:
        삭제된 디렉토리 수
    """
    import shutil
    import time

    if not STAGING_ROOT.exists():
        return 0

    cutoff = time.time() - (max_age_hours * 3600)
    removed = 0

    for d in STAGING_ROOT.iterdir():
        if d.is_dir():
            try:
                if d.stat().st_mtime < cutoff:
                    shutil.rmtree(d)
                    removed += 1
            except Exception:
                continue

    return removed