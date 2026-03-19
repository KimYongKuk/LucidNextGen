"""
PPT Generator MCP Server
사내 .pptx 템플릿 기반 프레젠테이션 생성 서버

템플릿 인덱싱 메타데이터를 참조하여 LLM이 구성한 slides JSON을 받아
실제 .pptx 파일을 생성한다.

지원 요소:
- TextBox (제목, 본문, 불릿 텍스트, 배경색/테두리 옵션)
- Table (사내 스타일 헤더, 교대행 색상, 셀 병합)
- Native Chart (line, column, bar, pie 등 + 시리즈 색상, 데이터라벨, 범례)
- Image (외부 이미지 삽입)
- CalloutBox (인사이트/경고/요약 강조 박스)
- KPICard (대시보드 핵심 수치 카드)
- Divider (시각 구분선)
"""

import sys
import os
import json
import re
import asyncio
from pathlib import Path
from datetime import datetime
from typing import Any

# backend/ 디렉토리를 sys.path에 추가 (MCP stdio 실행 시 필요)
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', '..')))

# MCP SDK
from mcp.server import Server
from mcp.server.stdio import stdio_server
from mcp.types import Tool, TextContent

# python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData, XyChartData

# 템플릿 인덱서
from app.mcp_servers.ppt_generator.template_indexer import (
    load_metadata, TEMPLATE_DIR, METADATA_FILENAME
)

# 경로 설정
OUTPUT_DIR = Path(__file__).resolve().parents[3] / "data" / "ppt_output"
CHART_OUTPUT_DIR = Path(__file__).resolve().parents[3] / "data" / "chart_output"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# MCP 서버 생성
server = Server("ppt-generator")

# ============================================================
# 메타데이터 캐시
# ============================================================
_cached_metadata: dict[str, Any] | None = None


def get_template_metadata() -> dict[str, Any] | None:
    """캐시된 메타데이터 반환"""
    global _cached_metadata
    if _cached_metadata is None:
        _cached_metadata = load_metadata()
    return _cached_metadata


def get_template_path() -> Path:
    """템플릿 .pptx 파일 경로 반환"""
    meta = get_template_metadata()
    if meta:
        return TEMPLATE_DIR / meta["template_file"]
    # 폴백: 기본 파일명
    return TEMPLATE_DIR / "PPT_Public.pptx"


# ============================================================
# 스타일 상수 (메타데이터에서 로드, 폴백 기본값)
# ============================================================

def _get_style() -> dict[str, Any]:
    """메타데이터의 style_guide 반환 (폴백 포함)"""
    meta = get_template_metadata()
    if meta and "style_guide" in meta:
        return meta["style_guide"]
    return {
        "font_family": "맑은 고딕",
        "theme_colors": {
            "accent1_blue": "4472C4",
            "accent2_orange": "ED7D31",
            "dk2_dark_gray": "44546A",
        },
        "table_style": {
            "header_fill": "182F54",
            "header_font_color": "FFFFFF",
            "header_font_size": 10,
            "header_bold": True,
            "body_font_size": 10,
            "body_alt_fill": "E7EAEE",
        },
    }


def _rgb(hex_color: str) -> RGBColor:
    """HEX 문자열 → RGBColor"""
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def _resolve_color(color_ref: str, style: dict) -> RGBColor:
    """색상 참조 해석 (HEX 직접 또는 테마 키)"""
    if not color_ref:
        return RGBColor(0, 0, 0)
    # 테마 색상 키 (예: "accent2_orange")
    tc = style.get("theme_colors", {})
    if color_ref in tc:
        return _rgb(tc[color_ref])
    # "white", "black" 등 이름
    named = {"white": "FFFFFF", "black": "000000", "red": "FF0000", "green": "00B050", "blue": "4472C4"}
    if color_ref.lower() in named:
        return _rgb(named[color_ref.lower()])
    # HEX 직접
    if len(color_ref) == 6:
        try:
            return _rgb(color_ref)
        except Exception:
            pass
    return RGBColor(0, 0, 0)


# ============================================================
# 유틸리티
# ============================================================

def sanitize_filename(filename: str) -> str:
    """파일명 정리 (특수문자 제거)"""
    return re.sub(r'[<>:"/\\|?*`\'"\s]', '_', filename).strip("_") or "presentation"


_DEFAULT_TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)  # 템플릿 테마 흰색 상속 방지


def _set_font(run, font_name: str, size_pt: int, bold: bool = False,
              color: RGBColor | None = None):
    """Run에 폰트 속성 설정 — color 미지정 시 기본 검정(333333) 적용"""
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color if color is not None else _DEFAULT_TEXT_COLOR


def _set_cell_font(cell, font_name: str, size_pt: int, bold: bool = False,
                   color: RGBColor | None = None, alignment: PP_ALIGN = PP_ALIGN.LEFT):
    """테이블 셀 텍스트에 폰트 적용"""
    for para in cell.text_frame.paragraphs:
        para.alignment = alignment
        for run in para.runs:
            _set_font(run, font_name, size_pt, bold, color)


# ============================================================
# TextBox 렌더링
# ============================================================

def render_textbox(slide, shape_def: dict, style: dict):
    """TextBox Shape 추가 (선택적 배경색/테두리 지원)"""
    left = Inches(shape_def.get("left", 0.37))
    top = Inches(shape_def.get("top", 1.15))
    width = Inches(shape_def.get("width", 12.0))
    height = Inches(shape_def.get("height", 0.5))

    fill_color = shape_def.get("fill_color", "")
    border_color = shape_def.get("border_color", "")

    # fill/border가 있으면 ROUNDED_RECTANGLE, 없으면 기존 textbox
    if fill_color or border_color:
        txBox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        # 모서리 반경 설정 (적당히 둥글게)
        txBox.adjustments[0] = 0.02
        if fill_color:
            txBox.fill.solid()
            txBox.fill.fore_color.rgb = _resolve_color(fill_color, style)
        else:
            txBox.fill.background()  # 투명
        if border_color:
            txBox.line.color.rgb = _resolve_color(border_color, style)
            txBox.line.width = Pt(shape_def.get("border_width", 1))
        else:
            txBox.line.fill.background()  # 테두리 없음
    else:
        txBox = slide.shapes.add_textbox(left, top, width, height)

    tf = txBox.text_frame
    tf.word_wrap = True

    # fill/border가 있을 때 내부 여백 설정
    if fill_color or border_color:
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.08)
        tf.margin_bottom = Inches(0.08)

    text = shape_def.get("text", "")
    font_size = shape_def.get("font_size", 10)
    bold = shape_def.get("bold", False)
    color_ref = shape_def.get("color", "")
    align = shape_def.get("alignment", "left")

    font_name = style.get("font_family", "맑은 고딕")
    color = _resolve_color(color_ref, style) if color_ref else None

    # 정렬 매핑
    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }
    pp_align = align_map.get(align, PP_ALIGN.LEFT)

    # 텍스트 줄 분리 (불릿 지원)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        para.alignment = pp_align

        # 불릿 처리
        stripped = line.lstrip()
        if stripped.startswith("- ") or stripped.startswith("• "):
            para.level = 0
            line_text = stripped[2:]
        elif stripped.startswith("  - ") or stripped.startswith("  • "):
            para.level = 1
            line_text = stripped[4:]
        else:
            line_text = line

        run = para.add_run()
        run.text = line_text
        _set_font(run, font_name, font_size, bold, color)

    return txBox


# ============================================================
# 테이블 렌더링
# ============================================================

def render_table(slide, shape_def: dict, style: dict):
    """테이블 Shape 추가 (사내 스타일 적용)"""
    table_data = shape_def.get("table", {})
    headers = table_data.get("headers", [])
    rows_data = table_data.get("rows", [])

    if not headers and not rows_data:
        return None

    # 위치/크기
    left = Inches(shape_def.get("left", 0.37))
    top = Inches(shape_def.get("top", 1.15))
    width = Inches(shape_def.get("width", 12.6))
    height = Inches(shape_def.get("height", 5.0))

    # 행/열 수 계산
    num_cols = len(headers) if headers else (len(rows_data[0]) if rows_data else 0)
    header_rows = table_data.get("header_rows", [])
    num_header_rows = len(header_rows) if header_rows else (1 if headers else 0)
    num_body_rows = len(rows_data)
    total_rows = num_header_rows + num_body_rows

    if total_rows == 0 or num_cols == 0:
        return None

    # 테이블 너비 보정: 최소 8인치, 기본 12.6인치
    table_width_inches = shape_def.get("width", 12.6)
    if table_width_inches < 8.0:
        table_width_inches = 12.6
        width = Inches(table_width_inches)

    # 테이블 생성
    table_shape = slide.shapes.add_table(total_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # 열 너비 설정: col_widths를 비율(weight)로 취급하여 전체 너비에 맞게 정규화
    col_widths = table_data.get("col_widths", [])
    if col_widths and len(col_widths) == num_cols:
        total_weight = sum(col_widths)
        if total_weight > 0:
            for i, w in enumerate(col_widths):
                normalized = (w / total_weight) * table_width_inches
                table.columns[i].width = Inches(normalized)
    else:
        # col_widths 미지정 시 균등 분배 (안전장치)
        equal_width = table_width_inches / num_cols
        for i in range(num_cols):
            table.columns[i].width = Inches(equal_width)

    # 스타일 정보
    ts = style.get("table_style", {})
    font_name = style.get("font_family", "맑은 고딕")
    header_fill = ts.get("header_fill", "182F54")
    header_font_color = ts.get("header_font_color", "FFFFFF")
    header_font_size = ts.get("header_font_size", 10)
    body_font_size = ts.get("body_font_size", 10)
    body_alt_fill = ts.get("body_alt_fill", "E7EAEE")

    # LLM이 스타일 오버라이드 가능
    header_fill = table_data.get("header_fill", header_fill)
    body_alt_fill_enabled = table_data.get("alt_row_fill", True)

    # 병합 헤더 (header_rows)
    if header_rows:
        for row_idx, hrow in enumerate(header_rows):
            for cell_def in hrow:
                col_start = cell_def.get("col", 0)
                colspan = cell_def.get("colspan", 1)
                rowspan = cell_def.get("rowspan", 1)
                text = cell_def.get("text", "")

                cell = table.cell(row_idx, col_start)
                cell.text = text

                # 셀 병합
                if colspan > 1 or rowspan > 1:
                    end_row = row_idx + rowspan - 1
                    end_col = col_start + colspan - 1
                    cell.merge(table.cell(end_row, end_col))

                # 스타일
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(header_fill)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                _set_cell_font(cell, font_name, header_font_size, True,
                               _rgb(header_font_color), PP_ALIGN.CENTER)
    elif headers:
        # 단순 헤더
        for col_idx, header_text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header_text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(header_fill)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            _set_cell_font(cell, font_name, header_font_size, True,
                           _rgb(header_font_color), PP_ALIGN.CENTER)

    # 본문 행
    for row_idx, row_data in enumerate(rows_data):
        actual_row = row_idx + num_header_rows
        for col_idx, cell_value in enumerate(row_data):
            if col_idx >= num_cols:
                break
            cell = table.cell(actual_row, col_idx)
            cell.text = str(cell_value) if cell_value is not None else ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 교대행 배경색
            if body_alt_fill_enabled and row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(body_alt_fill)

            _set_cell_font(cell, font_name, body_font_size, False,
                           RGBColor(0x33, 0x33, 0x33), PP_ALIGN.CENTER)

    # 본문 병합 (body_merges)
    body_merges = table_data.get("body_merges", [])
    for merge_def in body_merges:
        r1 = merge_def.get("row", 0) + num_header_rows
        c1 = merge_def.get("col", 0)
        rs = merge_def.get("rowspan", 1)
        cs = merge_def.get("colspan", 1)
        r2 = r1 + rs - 1
        c2 = c1 + cs - 1
        table.cell(r1, c1).merge(table.cell(r2, c2))

    return table_shape


# ============================================================
# CalloutBox 렌더링
# ============================================================

# 프리셋 스타일 (style 키 → fill, accent bar, icon)
CALLOUT_PRESETS = {
    "insight":  {"fill": "E8F0FE", "accent": "4472C4", "icon": "\U0001f4a1"},  # 💡
    "warning":  {"fill": "FFF3E0", "accent": "ED7D31", "icon": "\u26a0\ufe0f"},  # ⚠️
    "success":  {"fill": "E8F5E9", "accent": "70AD47", "icon": "\u2705"},       # ✅
    "summary":  {"fill": "F5F5F5", "accent": "44546A", "icon": "\U0001f4cb"},  # 📋
}


def render_callout_box(slide, shape_def: dict, style: dict):
    """강조 박스 (인사이트, 경고, 요약 등) — 좌측 accent bar + 아이콘 + 텍스트"""
    left = Inches(shape_def.get("left", 0.37))
    top = Inches(shape_def.get("top", 5.8))
    width = Inches(shape_def.get("width", 12.6))
    height = Inches(shape_def.get("height", 0.8))

    font_name = style.get("font_family", "맑은 고딕")
    preset_name = shape_def.get("style", "insight")
    preset = CALLOUT_PRESETS.get(preset_name, CALLOUT_PRESETS["insight"])

    fill_hex = shape_def.get("fill_color", preset["fill"])
    accent_hex = shape_def.get("accent_color", preset["accent"])
    # accent_color가 테마 키면 해석
    if accent_hex in style.get("theme_colors", {}):
        accent_hex = style["theme_colors"][accent_hex]
    icon = shape_def.get("icon", preset["icon"])
    text = shape_def.get("text", "")
    font_size = shape_def.get("font_size", 10)

    # 1) 배경 라운드 사각형
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    bg.adjustments[0] = 0.03
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(fill_hex)
    bg.line.fill.background()

    # 2) 좌측 accent bar (얇은 사각형)
    bar_width = Inches(0.06)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, bar_width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(accent_hex)
    bar.line.fill.background()

    # 3) 텍스트 (아이콘 + 내용)
    text_left = Inches(shape_def.get("left", 0.37)) + Inches(0.2)
    text_width = width - Inches(0.35)
    txBox = slide.shapes.add_textbox(text_left, top, text_width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)

    # 세로 중앙 정렬
    from lxml import etree
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    bodyPr = tf._txBody.find('.//a:bodyPr', nsmap)
    if bodyPr is not None:
        bodyPr.set('anchor', 'ctr')

    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    # 아이콘
    if icon:
        icon_run = para.add_run()
        icon_run.text = f"{icon} "
        _set_font(icon_run, font_name, font_size)
    # 텍스트
    run = para.add_run()
    run.text = text
    _set_font(run, font_name, font_size, color=_rgb("333333"))

    return bg


# ============================================================
# KPI Card 렌더링
# ============================================================

def render_kpi_card(slide, shape_def: dict, style: dict):
    """KPI 카드 — 큰 숫자 + 라벨 + 트렌드 표시"""
    left = Inches(shape_def.get("left", 0.37))
    top = Inches(shape_def.get("top", 1.15))
    width = Inches(shape_def.get("width", 2.9))
    height = Inches(shape_def.get("height", 1.8))

    font_name = style.get("font_family", "맑은 고딕")
    accent_ref = shape_def.get("accent_color", "accent1_blue")
    accent_rgb = _resolve_color(accent_ref, style)

    value = shape_def.get("value", "0")
    label = shape_def.get("label", "")
    trend = shape_def.get("trend", "")
    trend_direction = shape_def.get("trend_direction", "")  # up / down / neutral
    value_size = shape_def.get("value_size", 36)

    # 1) 배경 카드
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    bg.adjustments[0] = 0.05
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb("F8F9FA")
    bg.line.color.rgb = _rgb("E0E0E0")
    bg.line.width = Pt(0.75)

    # 2) 상단 accent line
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(4)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = accent_rgb
    accent_line.line.fill.background()

    # 3) 값 (큰 숫자)
    val_top = top + Inches(0.25)
    val_box = slide.shapes.add_textbox(left, val_top, width, Inches(0.7))
    val_tf = val_box.text_frame
    val_tf.word_wrap = True
    val_para = val_tf.paragraphs[0]
    val_para.alignment = PP_ALIGN.CENTER
    val_run = val_para.add_run()
    val_run.text = str(value)
    _set_font(val_run, font_name, value_size, bold=True, color=accent_rgb)

    # 4) 라벨
    if label:
        lbl_top = val_top + Inches(0.7)
        lbl_box = slide.shapes.add_textbox(left, lbl_top, width, Inches(0.35))
        lbl_tf = lbl_box.text_frame
        lbl_tf.word_wrap = True
        lbl_para = lbl_tf.paragraphs[0]
        lbl_para.alignment = PP_ALIGN.CENTER
        lbl_run = lbl_para.add_run()
        lbl_run.text = label
        _set_font(lbl_run, font_name, 10, color=_rgb("666666"))

    # 5) 트렌드
    if trend:
        trn_top = val_top + Inches(1.05)
        trn_box = slide.shapes.add_textbox(left, trn_top, width, Inches(0.3))
        trn_tf = trn_box.text_frame
        trn_para = trn_tf.paragraphs[0]
        trn_para.alignment = PP_ALIGN.CENTER

        arrow = ""
        trn_color = _rgb("666666")
        if trend_direction == "up":
            arrow = "▲ "
            trn_color = _rgb("00B050")  # 녹색
        elif trend_direction == "down":
            arrow = "▼ "
            trn_color = _rgb("FF0000")  # 빨강

        trn_run = trn_para.add_run()
        trn_run.text = f"{arrow}{trend}"
        _set_font(trn_run, font_name, 11, bold=True, color=trn_color)

    return bg


# ============================================================
# Divider (구분선) 렌더링
# ============================================================

def render_divider(slide, shape_def: dict, style: dict):
    """시각 구분선 — 얇은 사각형"""
    left = Inches(shape_def.get("left", 0.37))
    top = Inches(shape_def.get("top", 3.5))
    width = Inches(shape_def.get("width", 12.6))
    thickness = shape_def.get("thickness", 1.5)
    color_ref = shape_def.get("color", "E0E0E0")

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(thickness))
    line.fill.solid()
    line.fill.fore_color.rgb = _resolve_color(color_ref, style)
    line.line.fill.background()

    return line


# ============================================================
# 네이티브 차트 렌더링
# ============================================================

CHART_TYPE_MAP = {
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
    "pie": XL_CHART_TYPE.PIE,
    "area": XL_CHART_TYPE.AREA,
    "area_stacked": XL_CHART_TYPE.AREA_STACKED,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
}


def render_native_chart(slide, shape_def: dict, style: dict):
    """네이티브 PPT 차트 추가 (편집 가능, 스타일 옵션 지원)"""
    chart_info = shape_def.get("chart", {})
    chart_type_str = chart_info.get("chart_type", "column")
    categories = chart_info.get("categories", [])
    series_list = chart_info.get("series", [])

    if not categories or not series_list:
        return None

    xl_type = CHART_TYPE_MAP.get(chart_type_str)
    if xl_type is None:
        return None

    left = Inches(shape_def.get("left", 0.5))
    top = Inches(shape_def.get("top", 1.5))
    width = Inches(shape_def.get("width", 12.0))
    height = Inches(shape_def.get("height", 5.0))

    # scatter는 XyChartData 사용
    if chart_type_str == "scatter":
        chart_data = XyChartData()
        for series in series_list:
            xy_series = chart_data.add_series(series.get("name", ""))
            x_values = categories
            y_values = series.get("values", [])
            for x, y in zip(x_values, y_values):
                xy_series.add_data_point(x, y)
    else:
        chart_data = CategoryChartData()
        chart_data.categories = categories
        for series in series_list:
            chart_data.add_series(
                series.get("name", ""),
                series.get("values", [])
            )

    chart_shape = slide.shapes.add_chart(
        xl_type, left, top, width, height, chart_data
    )
    chart = chart_shape.chart
    font_name = style.get("font_family", "맑은 고딕")

    # 차트 제목 설정
    chart_title = chart_info.get("title", "")
    if chart_title:
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_title
        for para in chart.chart_title.text_frame.paragraphs:
            for run in para.runs:
                _set_font(run, font_name, 12, True)

    # --- 시리즈 색상 ---
    series_colors = chart_info.get("series_colors", [])
    if series_colors:
        is_pie = chart_type_str in ("pie", "doughnut")
        if is_pie and len(chart.series) > 0:
            # pie/doughnut: 포인트별 색상
            series_obj = chart.series[0]
            for i, pt in enumerate(series_obj.points):
                if i < len(series_colors):
                    pt.format.fill.solid()
                    pt.format.fill.fore_color.rgb = _rgb(series_colors[i])
        else:
            # 일반 차트: 시리즈별 색상
            for i, series_obj in enumerate(chart.series):
                if i < len(series_colors):
                    series_obj.format.fill.solid()
                    series_obj.format.fill.fore_color.rgb = _rgb(series_colors[i])
                    # line 계열은 선 색상도 설정
                    if chart_type_str in ("line", "line_markers"):
                        series_obj.format.line.color.rgb = _rgb(series_colors[i])

    # --- 데이터 라벨 ---
    if chart_info.get("data_labels", False):
        try:
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.font.size = Pt(9)
            data_labels.font.name = font_name
            num_fmt = chart_info.get("number_format", "")
            if num_fmt:
                data_labels.number_format = num_fmt
                data_labels.number_format_is_linked = False
        except Exception:
            pass  # 일부 차트 타입에서는 데이터 라벨 미지원

    # --- 범례 위치 ---
    legend_pos = chart_info.get("legend_position", "")
    if legend_pos == "none":
        chart.has_legend = False
    elif legend_pos:
        chart.has_legend = True
        pos_map = {
            "bottom": XL_LEGEND_POSITION.BOTTOM,
            "right": XL_LEGEND_POSITION.RIGHT,
            "top": XL_LEGEND_POSITION.TOP,
            "left": XL_LEGEND_POSITION.LEFT,
        }
        if legend_pos in pos_map:
            chart.legend.position = pos_map[legend_pos]
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(9)
            chart.legend.font.name = font_name

    return chart_shape


# ============================================================
# 이미지 렌더링
# ============================================================

def render_image(slide, shape_def: dict, style: dict):
    """이미지 Shape 추가 (차트 이미지 등)"""
    image_path = shape_def.get("path", "")
    if not image_path:
        return None

    # 상대 경로 → 절대 경로
    img_path = Path(image_path)
    if not img_path.is_absolute():
        img_path = CHART_OUTPUT_DIR / image_path

    if not img_path.exists():
        return None

    left = Inches(shape_def.get("left", 1.0))
    top = Inches(shape_def.get("top", 1.5))
    width = Inches(shape_def.get("width", 10.0))
    height = Inches(shape_def.get("height", 5.0))

    return slide.shapes.add_picture(str(img_path), left, top, width, height)


# ============================================================
# 내용 슬라이드 헤더 렌더링
# ============================================================

def render_content_header(slide, slide_def: dict, style: dict, doc_title: str = ""):
    """내용 레이아웃 슬라이드의 공통 헤더 요소 렌더링"""
    meta = get_template_metadata()
    content_area = {}
    if meta:
        content_area = meta.get("content_areas", {}).get("내용", {}).get("header", {})

    font_name = style.get("font_family", "맑은 고딕")
    tc = style.get("theme_colors", {})
    orange = _resolve_color("accent2_orange", style)

    # 1. 문서명 (좌측 상단)
    header_doc_title = slide_def.get("doc_title", doc_title)
    if header_doc_title:
        dt_area = content_area.get("doc_title", {"left": 0.37, "top": 0.30, "width": 3.0, "height": 0.12})
        txBox = slide.shapes.add_textbox(
            Inches(dt_area["left"]), Inches(dt_area["top"]),
            Inches(dt_area["width"]), Inches(dt_area["height"])
        )
        run = txBox.text_frame.paragraphs[0].add_run()
        run.text = header_doc_title
        _set_font(run, font_name, dt_area.get("font_size", 7))

    # 2. 목차 경로 (breadcrumb)
    breadcrumb = slide_def.get("breadcrumb", "")
    if breadcrumb:
        bc_area = content_area.get("breadcrumb", {"left": 0.33, "top": 0.44, "width": 2.0, "height": 0.15})
        txBox = slide.shapes.add_textbox(
            Inches(bc_area["left"]), Inches(bc_area["top"]),
            Inches(bc_area["width"]), Inches(bc_area["height"])
        )
        run = txBox.text_frame.paragraphs[0].add_run()
        run.text = breadcrumb
        _set_font(run, font_name, bc_area.get("font_size", 9), color=orange)

    # 3. 메인 제목
    main_title = slide_def.get("title", "")
    if main_title:
        mt_area = content_area.get("main_title", {"left": 0.37, "top": 0.69, "width": 12.0, "height": 0.42})
        txBox = slide.shapes.add_textbox(
            Inches(mt_area["left"]), Inches(mt_area["top"]),
            Inches(mt_area["width"]), Inches(mt_area["height"])
        )
        run = txBox.text_frame.paragraphs[0].add_run()
        run.text = main_title
        _set_font(run, font_name, mt_area.get("font_size", 25), bold=True)

    # 4. 부제 (subtitle)
    subtitle = slide_def.get("subtitle", "")
    if subtitle:
        txBox = slide.shapes.add_textbox(
            Inches(0.37), Inches(1.15), Inches(6.0), Inches(0.25)
        )
        run = txBox.text_frame.paragraphs[0].add_run()
        run.text = subtitle
        _set_font(run, font_name, 15)


# ============================================================
# 슬라이드 렌더링 (메인)
# ============================================================

def render_slide(prs, slide_def: dict, style: dict, doc_title: str = ""):
    """단일 슬라이드 렌더링"""
    layout_index = slide_def.get("layout_index", 3)  # 기본: 내용
    layout_name = slide_def.get("layout_name", "")

    # 디버그 로깅: LLM이 보내는 슬라이드 데이터 확인
    slide_keys = list(slide_def.keys())
    print(f"[PPT] render_slide: layout_index={layout_index}, layout_name={layout_name}, keys={slide_keys}")

    # 레이아웃 이름으로 인덱스 찾기
    if layout_name and not slide_def.get("layout_index"):
        meta = get_template_metadata()
        if meta:
            for l in meta["layouts"]:
                if l["name"] == layout_name:
                    layout_index = l["index"]
                    break

    # 범위 체크
    if layout_index >= len(prs.slide_layouts):
        layout_index = min(3, len(prs.slide_layouts) - 1)

    layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)

    # 레이아웃별 처리
    layout_actual_name = layout.name

    if layout_actual_name == "내용":
        # 공통 헤더 렌더링
        render_content_header(slide, slide_def, style, doc_title)

        # shapes 배열의 각 요소 렌더링
        for shape_def in slide_def.get("shapes", []):
            shape_type = shape_def.get("type", "")
            if shape_type == "textbox":
                render_textbox(slide, shape_def, style)
            elif shape_type == "table":
                render_table(slide, shape_def, style)
            elif shape_type == "chart":
                render_native_chart(slide, shape_def, style)
            elif shape_type == "image":
                render_image(slide, shape_def, style)
            elif shape_type == "callout_box":
                render_callout_box(slide, shape_def, style)
            elif shape_type == "kpi_card":
                render_kpi_card(slide, shape_def, style)
            elif shape_type == "divider":
                render_divider(slide, shape_def, style)

    elif layout_actual_name == "표지":
        # 표지 슬라이드 콘텐츠
        font_name = style.get("font_family", "맑은 고딕")
        meta = get_template_metadata()
        cover_area = {}
        if meta:
            cover_area = meta.get("content_areas", {}).get("표지", {})

        title = slide_def.get("title", "")
        if title:
            area = cover_area.get("title_area", {"left": 1.67, "top": 2.83, "width": 10.0, "height": 0.9})
            txBox = slide.shapes.add_textbox(
                Inches(area["left"]), Inches(area["top"]),
                Inches(area["width"]), Inches(area["height"])
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.auto_size = None
            # 세로 중앙 정렬 (MSO_ANCHOR.MIDDLE)
            from lxml import etree
            nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            bodyPr = tf._txBody.find('.//a:bodyPr', nsmap)
            if bodyPr is not None:
                bodyPr.set('anchor', 'ctr')
            run = tf.paragraphs[0].add_run()
            run.text = title
            _set_font(run, font_name, 40, bold=True)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        date_text = slide_def.get("date", "")
        if date_text:
            area = cover_area.get("date_area", {"left": 4.5, "top": 5.0, "width": 4.33, "height": 0.35})
            txBox = slide.shapes.add_textbox(
                Inches(area["left"]), Inches(area["top"]),
                Inches(area["width"]), Inches(area["height"])
            )
            run = txBox.text_frame.paragraphs[0].add_run()
            run.text = date_text
            _set_font(run, font_name, 14)
            txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        author = slide_def.get("author", "")
        if author:
            area = cover_area.get("author_area", {"left": 4.5, "top": 5.6, "width": 4.33, "height": 0.35})
            txBox = slide.shapes.add_textbox(
                Inches(area["left"]), Inches(area["top"]),
                Inches(area["width"]), Inches(area["height"])
            )
            run = txBox.text_frame.paragraphs[0].add_run()
            run.text = author
            _set_font(run, font_name, 11)
            txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    elif layout_actual_name == "목차":
        font_name = style.get("font_family", "맑은 고딕")
        orange = _resolve_color("accent2_orange", style)

        # CONTENTS 타이틀
        title = slide_def.get("title", "CONTENTS")
        txBox = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.0), Inches(3.0), Inches(0.5)
        )
        txBox.text_frame.word_wrap = True
        run = txBox.text_frame.paragraphs[0].add_run()
        run.text = title
        _set_font(run, font_name, 30, bold=True)

        # 목차 항목 (items, sections, contents 등 다양한 키 지원)
        items = slide_def.get("items", [])
        if not items:
            items = slide_def.get("sections", [])
        if not items:
            items = slide_def.get("contents", [])
        if not items:
            items = slide_def.get("toc_items", [])

        print(f"[PPT] 목차 items: count={len(items)}, data={items[:3] if items else 'EMPTY'}")

        y_pos = 2.0
        for item in items:
            if isinstance(item, str):
                # 단순 문자열
                txBox = slide.shapes.add_textbox(
                    Inches(0.75), Inches(y_pos), Inches(7.0), Inches(0.35)
                )
                txBox.text_frame.word_wrap = True
                run = txBox.text_frame.paragraphs[0].add_run()
                run.text = item
                _set_font(run, font_name, 15, bold=True, color=orange)
                y_pos += 0.5
            elif isinstance(item, dict):
                # 대목차 (major, title, text, section, name 등 다양한 키 지원)
                major = (item.get("major", "")
                         or item.get("title", "")
                         or item.get("text", "")
                         or item.get("section", "")
                         or item.get("name", "")
                         or item.get("label", ""))
                if major:
                    txBox = slide.shapes.add_textbox(
                        Inches(0.75), Inches(y_pos), Inches(7.0), Inches(0.35)
                    )
                    txBox.text_frame.word_wrap = True
                    run = txBox.text_frame.paragraphs[0].add_run()
                    run.text = str(major)
                    _set_font(run, font_name, 15, bold=True, color=orange)
                    y_pos += 0.4

                # 소목차 (minor, sub_items, children, subsections 등 다양한 키 지원)
                minor_items = (item.get("minor", [])
                               or item.get("sub_items", [])
                               or item.get("children", [])
                               or item.get("subsections", [])
                               or item.get("items", []))
                for minor in minor_items:
                    minor_text = minor if isinstance(minor, str) else str(minor.get("title", minor.get("text", minor.get("name", str(minor)))))
                    txBox = slide.shapes.add_textbox(
                        Inches(1.5), Inches(y_pos), Inches(6.0), Inches(0.3)
                    )
                    txBox.text_frame.word_wrap = True
                    run = txBox.text_frame.paragraphs[0].add_run()
                    run.text = minor_text
                    _set_font(run, font_name, 15)
                    y_pos += 0.35

                y_pos += 0.15

    elif layout_actual_name == "간지":
        font_name = style.get("font_family", "맑은 고딕")
        white = RGBColor(0xFF, 0xFF, 0xFF)

        title = slide_def.get("title", "")
        if title:
            txBox = slide.shapes.add_textbox(
                Inches(0.7), Inches(2.5), Inches(6.5), Inches(1.1)
            )
            run = txBox.text_frame.paragraphs[0].add_run()
            run.text = title
            _set_font(run, font_name, 40, bold=True, color=white)

        subtitle = slide_def.get("subtitle", "")
        if subtitle:
            txBox = slide.shapes.add_textbox(
                Inches(0.7), Inches(3.8), Inches(6.5), Inches(0.5)
            )
            run = txBox.text_frame.paragraphs[0].add_run()
            run.text = subtitle
            _set_font(run, font_name, 20, color=white)

    elif layout_actual_name == "E.O.D":
        # E.O.D 레이아웃은 이미 "End Of Document" 텍스트가 있으므로
        # 추가 콘텐츠가 필요한 경우만 처리
        custom_text = slide_def.get("title", "")
        if custom_text:
            font_name = style.get("font_family", "맑은 고딕")
            txBox = slide.shapes.add_textbox(
                Inches(4.0), Inches(4.5), Inches(5.0), Inches(0.5)
            )
            run = txBox.text_frame.paragraphs[0].add_run()
            run.text = custom_text
            _set_font(run, font_name, 14)
            txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


# ============================================================
# 프레젠테이션 생성 (핵심)
# ============================================================

def create_pptx(slides: list[dict], filename: str, template: str = "",
                metadata: dict | None = None) -> tuple[Path, int]:
    """
    템플릿 기반 .pptx 파일 생성

    Args:
        slides: 슬라이드 정의 배열
        filename: 출력 파일명 (확장자 제외)
        template: 템플릿 이름 (미사용 시 기본 템플릿)
        metadata: 문서 메타데이터 (author, department, date 등)

    Returns:
        (output_path, slide_count)
    """
    # 템플릿 로드
    template_path = get_template_path()
    prs = Presentation(str(template_path))

    # 기존 슬라이드 모두 제거 (템플릿의 예시 슬라이드)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # 템플릿 섹션 목록 제거 (좌측 패널에서 모든 슬라이드가 "표지" 섹션 아래 표시되는 문제 방지)
    from lxml import etree
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    prs_xml = prs.element
    for section_lst in prs_xml.findall('.//p:sectionLst', nsmap):
        section_lst.getparent().remove(section_lst)
    # extLst 내부의 sectionLst도 제거 (p14 namespace)
    ns_p14 = 'http://schemas.microsoft.com/office/powerpoint/2010/main'
    for section_lst in prs_xml.findall(f'.//{{{ns_p14}}}sectionLst'):
        section_lst.getparent().remove(section_lst)

    style = _get_style()

    # 문서 제목 (헤더용)
    doc_title = ""
    if metadata:
        doc_title = metadata.get("title", "")
    if not doc_title and slides:
        # 첫 번째 표지 슬라이드의 제목 사용
        for s in slides:
            if s.get("layout_index", 3) == 0 or s.get("layout_name") == "표지":
                doc_title = s.get("title", "")
                break

    # 각 슬라이드 렌더링
    for slide_def in slides:
        render_slide(prs, slide_def, style, doc_title)

    # 파일 저장
    safe_filename = sanitize_filename(filename)
    output_path = OUTPUT_DIR / f"{safe_filename}.pptx"
    prs.save(str(output_path))

    return output_path, len(slides)


# ============================================================
# MCP Tool 정의
# ============================================================

@server.list_tools()
async def list_tools():
    """사용 가능한 도구 목록"""
    return [
        Tool(
            name="create_presentation",
            description="""사내 템플릿 기반 PPT 프레젠테이션을 생성합니다.

템플릿 레이아웃:
- [0] 표지: 제목, 날짜, 작성자
- [1] 목차: 대목차/소목차 나열
- [2] 간지: 섹션 구분 (대목차 강조)
- [3] 내용: 메인 콘텐츠 (텍스트, 테이블, 차트, 이미지 등 자유 배치)
- [4] E.O.D: 끝 페이지

내용 슬라이드 헤더 요소:
- doc_title: 좌측 상단 문서명 (7pt)
- breadcrumb: 목차 경로 (9pt, Orange)
- title: 메인 제목 (25pt, Bold)
- subtitle: 부제 (15pt)

내용 슬라이드 shapes 타입 (7종):
- textbox: 텍스트 (left, top, width, height, text, font_size, bold, color, alignment, fill_color?, border_color?, border_width?)
- table: 테이블 (left, top, width=12.6, height, table: {headers, rows, header_rows, col_widths(비율!), body_merges})
- chart: 네이티브 차트 (left, top, width, height, chart: {chart_type, categories, series, title, series_colors?, data_labels?, legend_position?, number_format?})
- image: 이미지 (left, top, width, height, path)
- callout_box: 강조 박스 (left, top, width, height, text, style=insight|warning|success|summary, icon?, fill_color?, accent_color?, font_size?)
- kpi_card: 수치 카드 (left, top, width, height, value, label, trend?, trend_direction=up|down?, accent_color?, value_size?)
- divider: 구분선 (left, top, width, color?, thickness?)

차트 타입: line, line_markers, column, column_stacked, bar, bar_stacked, pie, area, area_stacked, scatter, doughnut
차트 series_colors: HEX 배열 (예: ["4472C4","ED7D31"]). pie/doughnut은 포인트별 색상.
차트 legend_position: bottom, right, top, left, none

테이블 병합 헤더 (header_rows):
  [{"col": 0, "colspan": 2, "text": "그룹 헤더"}, ...]

⚠️ 테이블 크기 필수: width는 12.0~12.6 사용! col_widths는 비율(weight)로 지정 (예: [1,3] = 1:3 비율)

본문 영역: L=0.37, T=1.15, W=12.6, H=5.95 (inches)""",
            inputSchema={
                "type": "object",
                "properties": {
                    "slides": {
                        "type": "array",
                        "description": "슬라이드 정의 배열",
                        "items": {
                            "type": "object",
                            "properties": {
                                "layout_index": {
                                    "type": "integer",
                                    "description": "레이아웃 인덱스 (0=표지, 1=목차, 2=간지, 3=내용, 4=E.O.D)"
                                },
                                "layout_name": {
                                    "type": "string",
                                    "description": "레이아웃 이름 (표지, 목차, 간지, 내용, E.O.D)"
                                },
                                "title": {"type": "string"},
                                "subtitle": {"type": "string"},
                                "doc_title": {"type": "string"},
                                "breadcrumb": {"type": "string"},
                                "date": {"type": "string"},
                                "author": {"type": "string"},
                                "items": {
                                    "type": "array",
                                    "description": "목차 항목 (목차 레이아웃용). 문자열 배열 [\"I. 개요\", \"II. 현황\"] 또는 대목차/소목차 [{\"major\": \"I. 개요\", \"minor\": [\"배경\", \"목적\"]}]",
                                    "items": {
                                        "oneOf": [
                                            {"type": "string"},
                                            {
                                                "type": "object",
                                                "properties": {
                                                    "major": {"type": "string", "description": "대목차 텍스트"},
                                                    "minor": {"type": "array", "items": {"type": "string"}, "description": "소목차 텍스트 배열"}
                                                }
                                            }
                                        ]
                                    }
                                },
                                "shapes": {
                                    "type": "array",
                                    "description": "콘텐츠 Shape 배열 (내용 레이아웃용)",
                                    "items": {
                                        "type": "object",
                                        "properties": {
                                            "type": {
                                                "type": "string",
                                                "enum": ["textbox", "table", "chart", "image", "callout_box", "kpi_card", "divider"]
                                            },
                                            "left": {"type": "number"},
                                            "top": {"type": "number"},
                                            "width": {"type": "number"},
                                            "height": {"type": "number"},
                                        }
                                    }
                                }
                            }
                        }
                    },
                    "filename": {
                        "type": "string",
                        "description": "출력 파일명 (확장자 제외)"
                    },
                    "template": {
                        "type": "string",
                        "description": "템플릿 이름 (기본: PPT_Public)",
                        "default": "PPT_Public"
                    },
                    "metadata": {
                        "type": "object",
                        "description": "문서 메타데이터",
                        "properties": {
                            "title": {"type": "string"},
                            "author": {"type": "string"},
                            "department": {"type": "string"},
                            "date": {"type": "string"}
                        }
                    }
                },
                "required": ["slides", "filename"]
            }
        ),
        Tool(
            name="list_ppt_templates",
            description="사용 가능한 PPT 템플릿의 레이아웃 메타데이터를 반환합니다. 레이아웃 인덱스, 이름, 콘텐츠 영역, 스타일 가이드 등을 확인할 수 있습니다.",
            inputSchema={
                "type": "object",
                "properties": {},
            }
        ),
        Tool(
            name="list_generated_ppts",
            description="생성된 PPT 파일 목록을 반환합니다.",
            inputSchema={
                "type": "object",
                "properties": {},
            }
        ),
    ]


@server.call_tool()
async def call_tool(name: str, arguments: dict):
    """도구 실행"""

    if name == "create_presentation":
        slides = arguments.get("slides", [])
        filename = arguments.get("filename", "presentation")
        template = arguments.get("template", "")
        metadata = arguments.get("metadata", {})

        if not slides:
            return [TextContent(type="text", text="슬라이드 데이터가 비어있습니다.")]

        try:
            output_path, slide_count = create_pptx(slides, filename, template, metadata)
            file_size = output_path.stat().st_size / 1024

            return [TextContent(
                type="text",
                text=(
                    f"PPT 생성 완료\n\n"
                    f"파일명: {output_path.name}\n"
                    f"슬라이드: {slide_count}장\n"
                    f"크기: {file_size:.1f} KB\n"
                    f"경로: {output_path}"
                )
            )]

        except Exception as e:
            import traceback
            error_detail = traceback.format_exc()
            print(f"[PPT ERROR] {error_detail}")
            return [TextContent(
                type="text",
                text=(
                    f"PPT 생성 실패: {str(e)}\n\n"
                    f"⚠️ 재시도하지 말고, 현재까지 작성한 내용을 텍스트로 사용자에게 안내하세요.\n"
                    f"사용자에게 오류가 발생했음을 알리고, 슬라이드 구성을 간소화하여 다시 시도할 수 있음을 안내하세요.\n"
                    f"특히 table의 header_rows, body_merges 등 복잡한 셀 병합을 제거하고 단순 headers+rows로 재시도하세요."
                )
            )]

    elif name == "list_ppt_templates":
        try:
            meta = get_template_metadata()
            if not meta:
                return [TextContent(type="text", text="템플릿 메타데이터가 없습니다. template_indexer를 실행하세요.")]

            # LLM이 참고할 수 있는 요약 정보 반환
            from app.mcp_servers.ppt_generator.template_indexer import format_metadata_for_llm
            summary = format_metadata_for_llm(meta)

            return [TextContent(type="text", text=summary)]

        except Exception as e:
            return [TextContent(type="text", text=f"메타데이터 로드 실패: {str(e)}")]

    elif name == "list_generated_ppts":
        try:
            ppt_files = list(OUTPUT_DIR.glob("*.pptx"))

            if not ppt_files:
                return [TextContent(type="text", text="생성된 PPT 파일이 없습니다.")]

            file_list = []
            for f in sorted(ppt_files, key=lambda x: x.stat().st_mtime, reverse=True):
                stat = f.stat()
                size_kb = stat.st_size / 1024
                mtime = datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M')
                file_list.append(f"- {f.name} ({size_kb:.1f} KB) - {mtime}")

            return [TextContent(
                type="text",
                text=f"생성된 PPT 목록 ({len(ppt_files)}개)\n\n" + "\n".join(file_list)
            )]

        except Exception as e:
            return [TextContent(type="text", text=f"목록 조회 실패: {str(e)}")]

    return [TextContent(type="text", text=f"알 수 없는 도구: {name}")]


# ============================================================
# 서버 실행
# ============================================================

async def main():
    """MCP 서버 실행"""
    async with stdio_server() as (read_stream, write_stream):
        await server.run(
            read_stream,
            write_stream,
            server.create_initialization_options()
        )


if __name__ == "__main__":
    asyncio.run(main())
