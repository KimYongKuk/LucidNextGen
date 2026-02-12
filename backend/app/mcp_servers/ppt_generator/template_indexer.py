"""
PPT Template Indexer
.pptx 템플릿 파일을 파싱하여 LLM이 이해할 수 있는 레이아웃 메타데이터를 생성한다.

템플릿의 레이아웃, Shape 구성, 테마 색상, 폰트 스타일 등을 자동 추출하여
template_metadata.json 파일로 저장한다. 서버 런타임에서는 이 JSON만 로드하여 사용.

사용법:
    python -m app.mcp_servers.ppt_generator.template_indexer
    python -m app.mcp_servers.ppt_generator.template_indexer --template other.pptx

결과: template_metadata.json 파일이 템플릿과 같은 디렉토리에 생성됨
"""

import json
import sys
import zipfile
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# 기본 경로
TEMPLATE_DIR = Path(__file__).resolve().parents[3] / "data" / "ppt_template"
DEFAULT_TEMPLATE = "PPT_Public.pptx"
METADATA_FILENAME = "template_metadata.json"


def emu_to_inches(emu_value: int) -> float:
    """EMU(English Metric Units)를 inches로 변환"""
    if emu_value is None:
        return 0.0
    return round(emu_value / 914400, 2)


def _safe_text(shape) -> str:
    """Shape에서 안전하게 텍스트 추출"""
    try:
        if hasattr(shape, "text") and shape.text:
            return shape.text.replace("\n", " | ").strip()[:150]
    except Exception:
        pass
    return ""


def _shape_type_name(shape_type) -> str:
    """Shape type을 읽기 쉬운 문자열로 변환"""
    type_map = {
        MSO_SHAPE_TYPE.AUTO_SHAPE: "AUTO_SHAPE",
        MSO_SHAPE_TYPE.GROUP: "GROUP",
        MSO_SHAPE_TYPE.LINE: "LINE",
        MSO_SHAPE_TYPE.PICTURE: "PICTURE",
        MSO_SHAPE_TYPE.TEXT_BOX: "TEXT_BOX",
        MSO_SHAPE_TYPE.TABLE: "TABLE",
        MSO_SHAPE_TYPE.CHART: "CHART",
        MSO_SHAPE_TYPE.FREEFORM: "FREEFORM",
        MSO_SHAPE_TYPE.PLACEHOLDER: "PLACEHOLDER",
    }
    return type_map.get(shape_type, str(shape_type))


def _extract_theme_colors(pptx_path: Path) -> dict[str, str]:
    """ZIP에서 테마 색상 추출 (첫 번째 theme 파일)"""
    colors = {}
    try:
        with zipfile.ZipFile(str(pptx_path)) as z:
            for name in z.namelist():
                if "theme" in name.lower() and name.endswith(".xml"):
                    xml = z.read(name)
                    root = etree.fromstring(xml)
                    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
                    for scheme in root.iter(f"{{{ns_a}}}clrScheme"):
                        for child in scheme:
                            tag = child.tag.split("}")[-1]
                            for sub in child:
                                val = sub.attrib.get("val", sub.attrib.get("lastClr", ""))
                                if val:
                                    colors[tag] = val
                        break  # 첫 번째 color scheme만
                    break  # 첫 번째 theme 파일만
    except Exception:
        pass
    return colors


def _extract_layout_shapes(layout) -> list[dict[str, Any]]:
    """레이아웃의 Shape 목록 추출 (장식 요소 포함)"""
    shapes = []
    for shape in layout.shapes:
        shape_info: dict[str, Any] = {
            "type": _shape_type_name(shape.shape_type),
            "name": shape.name,
            "left": emu_to_inches(shape.left),
            "top": emu_to_inches(shape.top),
            "width": emu_to_inches(shape.width),
            "height": emu_to_inches(shape.height),
        }
        text = _safe_text(shape)
        if text:
            shape_info["text"] = text
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            shape_info["sub_shape_count"] = len(shape.shapes)
        shapes.append(shape_info)
    return shapes


def _analyze_content_slides(prs) -> dict[str, Any]:
    """기존 슬라이드를 분석하여 콘텐츠 패턴 추출"""
    content_analysis = {
        "header_elements": [],
        "table_positions": [],
        "total_slides": len(prs.slides),
    }

    for slide in prs.slides:
        if slide.slide_layout.name != "내용":
            continue

        for shape in slide.shapes:
            # 테이블 위치 수집
            if hasattr(shape, "has_table") and shape.has_table:
                content_analysis["table_positions"].append({
                    "left": emu_to_inches(shape.left),
                    "top": emu_to_inches(shape.top),
                    "width": emu_to_inches(shape.width),
                    "height": emu_to_inches(shape.height),
                    "rows": len(shape.table.rows),
                    "cols": len(shape.table.columns),
                })

    return content_analysis


def _extract_table_styles(prs) -> dict[str, Any]:
    """기존 슬라이드의 테이블에서 스타일 정보 추출"""
    header_fills = []
    body_fills = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not (hasattr(shape, "has_table") and shape.has_table):
                continue
            tbl = shape.table
            if len(tbl.rows) == 0:
                continue

            # 헤더 행 fill 색상
            for cell in tbl.rows[0].cells:
                try:
                    if cell.fill.type is not None:
                        try:
                            header_fills.append(str(cell.fill.fore_color.rgb))
                        except Exception:
                            pass
                except Exception:
                    pass

            # 본문 행 fill 색상
            if len(tbl.rows) > 1:
                for cell in tbl.rows[1].cells:
                    try:
                        if cell.fill.type is not None:
                            try:
                                body_fills.append(str(cell.fill.fore_color.rgb))
                            except Exception:
                                pass
                    except Exception:
                        pass

    # 가장 빈번한 색상
    result = {}
    if header_fills:
        from collections import Counter
        hc = Counter(header_fills)
        result["header_fill_color"] = hc.most_common(1)[0][0] if hc else None
    if body_fills:
        from collections import Counter
        bc = Counter(body_fills)
        result["body_alt_fill_color"] = bc.most_common(1)[0][0] if bc else None

    return result


def index_template(pptx_path: Path) -> dict[str, Any]:
    """
    .pptx 템플릿을 파싱하여 LLM이 이해할 수 있는 종합 메타데이터를 생성한다.
    """
    prs = Presentation(str(pptx_path))

    # 1. 테마 색상 추출
    theme_colors = _extract_theme_colors(pptx_path)

    # 2. 레이아웃 정보 추출
    layouts = []
    for idx, layout in enumerate(prs.slide_layouts):
        layout_info: dict[str, Any] = {
            "index": idx,
            "name": layout.name,
            "shapes": _extract_layout_shapes(layout),
            "placeholders": [],
        }

        # Placeholder 정보 (있는 경우)
        for ph in layout.placeholders:
            ph_info = {
                "idx": ph.placeholder_format.idx,
                "type": str(ph.placeholder_format.type),
                "name": ph.name,
                "left": emu_to_inches(ph.left),
                "top": emu_to_inches(ph.top),
                "width": emu_to_inches(ph.width),
                "height": emu_to_inches(ph.height),
            }
            layout_info["placeholders"].append(ph_info)

        layouts.append(layout_info)

    # 3. 기존 슬라이드 분석
    content_analysis = _analyze_content_slides(prs)
    table_styles = _extract_table_styles(prs)

    # 4. 콘텐츠 영역 정보 (기존 슬라이드 분석 기반)
    # "내용" 레이아웃의 콘텐츠 영역은 기존 슬라이드 패턴에서 추론
    content_areas = {}
    for layout in layouts:
        name = layout["name"]
        if name == "표지":
            content_areas[name] = {
                "description": "표지 슬라이드 - 제목, 날짜, 소속 배치 영역",
                "title_area": {"left": 1.67, "top": 2.83, "width": 10.0, "height": 0.9},
                "date_area": {"left": 4.5, "top": 5.0, "width": 4.33, "height": 0.35},
                "author_area": {"left": 4.5, "top": 5.6, "width": 4.33, "height": 0.35},
            }
        elif name == "목차":
            content_areas[name] = {
                "description": "목차 슬라이드 - 대목차/소목차 나열",
                "title_area": {"left": 0.75, "top": 1.0, "width": 3.0, "height": 0.5},
                "items_area": {"left": 0.75, "top": 2.0, "width": 5.5, "height": 4.5},
            }
        elif name == "간지":
            content_areas[name] = {
                "description": "섹션 간지 슬라이드 - 대목차 강조 표시",
                "title_area": {"left": 0.7, "top": 2.5, "width": 6.5, "height": 1.1},
                "subtitle_area": {"left": 0.7, "top": 3.8, "width": 6.5, "height": 2.5},
            }
        elif name == "내용":
            content_areas[name] = {
                "description": "메인 콘텐츠 슬라이드 - 자유 배치",
                "header": {
                    "doc_title": {"left": 0.37, "top": 0.30, "width": 3.0, "height": 0.12,
                                  "font_size": 7, "description": "문서 제목 (표지 타이틀)"},
                    "breadcrumb": {"left": 0.33, "top": 0.44, "width": 2.0, "height": 0.15,
                                   "font_size": 9, "color": "accent2_orange",
                                   "description": "목차 경로 (예: Ⅲ. 공정기술 전략 회의)"},
                    "main_title": {"left": 0.37, "top": 0.69, "width": 12.0, "height": 0.42,
                                   "font_size": 25, "bold": True,
                                   "description": "슬라이드 메인 제목"},
                },
                "body_area": {"left": 0.37, "top": 1.15, "width": 12.6, "height": 5.95,
                              "description": "본문 영역 (텍스트, 테이블, 차트 등)"},
                "footer_line_y": 7.21,
            }
        elif name == "E.O.D":
            content_areas[name] = {
                "description": "끝 슬라이드 - End Of Document",
                "title_area": {"left": 4.0, "top": 3.0, "width": 5.0, "height": 1.0},
            }

    # 5. 스타일 가이드 조합
    style_guide = {
        "font_family": "맑은 고딕",
        "theme_colors": {
            "dk1_text": theme_colors.get("dk1", "000000"),
            "lt1_background": theme_colors.get("lt1", "FFFFFF"),
            "dk2_dark_gray": theme_colors.get("dk2", "44546A"),
            "lt2_light_gray": theme_colors.get("lt2", "E7E6E6"),
            "accent1_blue": theme_colors.get("accent1", "4472C4"),
            "accent2_orange": theme_colors.get("accent2", "ED7D31"),
            "accent3_gray": theme_colors.get("accent3", "A5A5A5"),
            "accent4_gold": theme_colors.get("accent4", "FFC000"),
            "accent5_light_blue": theme_colors.get("accent5", "5B9BD5"),
            "accent6_green": theme_colors.get("accent6", "70AD47"),
        },
        "font_sizes": {
            "cover_title": {"pt": 40, "bold": True, "usage": "표지 보고서명"},
            "cover_date": {"pt": 14, "usage": "표지 날짜"},
            "cover_author": {"pt": 11, "usage": "표지 소속/작성자"},
            "toc_heading": {"pt": 30, "bold": True, "usage": "목차 CONTENTS"},
            "toc_major": {"pt": 15, "bold": True, "color": "accent2_orange", "usage": "목차 대분류"},
            "toc_minor": {"pt": 15, "usage": "목차 소분류"},
            "section_title": {"pt": 40, "bold": True, "color": "white", "usage": "간지 대목차"},
            "section_subtitle": {"pt": 20, "color": "white", "usage": "간지 소목차"},
            "content_doc_title": {"pt": 7, "usage": "내용 슬라이드 상단 문서명"},
            "content_breadcrumb": {"pt": 9, "color": "accent2_orange", "usage": "내용 슬라이드 목차 경로"},
            "content_main_title": {"pt": 25, "bold": True, "usage": "내용 슬라이드 메인 제목"},
            "content_subtitle": {"pt": 15, "usage": "내용 슬라이드 부제/헤드 메시지"},
            "content_section_title": {"pt": 12, "bold": True, "color": "accent2_orange", "usage": "내용 섹션 타이틀"},
            "content_body": {"pt": 10, "usage": "본문 텍스트"},
            "content_body_large": {"pt": 12, "bold": True, "usage": "본문 강조 텍스트"},
            "eod_title": {"pt": 20, "bold": True, "usage": "마무리 텍스트"},
        },
        "table_style": {
            "header_fill": table_styles.get("header_fill_color", "4472C4"),
            "header_font_color": "FFFFFF",
            "header_font_size": 10,
            "header_bold": True,
            "body_font_size": 10,
            "body_alt_fill": table_styles.get("body_alt_fill_color", "E7EAEE"),
            "dark_navy_header": "182F54",
        },
        "footer_tab": {
            "description": "내용 슬라이드 우측 상단에 문서 유형 표시 탭",
            "font_size": 11,
            "bold": True,
            "color": "white",
            "background": "accent1_blue",
        },
    }

    # 6. 최종 메타데이터 조합
    metadata: dict[str, Any] = {
        "template_name": pptx_path.stem,
        "template_file": pptx_path.name,
        "slide_width_inches": emu_to_inches(prs.slide_width),
        "slide_height_inches": emu_to_inches(prs.slide_height),
        "layouts": layouts,
        "content_areas": content_areas,
        "style_guide": style_guide,
        "analysis_notes": {
            "placeholder_based": False,
            "description": (
                "이 템플릿은 placeholder가 아닌 일반 Shape으로 구성됨. "
                "레이아웃은 배경 장식(로고, 선, 페이지 번호 등)을 제공하고, "
                "콘텐츠는 새로운 Shape(TextBox, Table, Chart, Image)을 "
                "content_areas 영역에 추가하는 방식으로 생성."
            ),
            "primary_content_layout": "내용 (index 3)",
            "total_existing_slides": content_analysis["total_slides"],
            "table_count_in_existing": len(content_analysis["table_positions"]),
        },
    }

    return metadata


def save_metadata(metadata: dict[str, Any], output_path: Path) -> None:
    """메타데이터를 JSON 파일로 저장"""
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(metadata, f, ensure_ascii=False, indent=2)


def load_metadata(metadata_path: Path | None = None) -> dict[str, Any] | None:
    """저장된 메타데이터 JSON 파일 로드 (서버 런타임용)"""
    if metadata_path is None:
        metadata_path = TEMPLATE_DIR / METADATA_FILENAME
    if not metadata_path.exists():
        return None
    with open(metadata_path, "r", encoding="utf-8") as f:
        return json.load(f)


def format_metadata_for_llm(metadata: dict[str, Any]) -> str:
    """메타데이터를 LLM 시스템 프롬프트용 텍스트로 변환"""
    lines = []
    lines.append(f"## 템플릿: {metadata['template_name']}")
    lines.append(f"슬라이드 크기: {metadata['slide_width_inches']}\" x {metadata['slide_height_inches']}\"")
    lines.append("")

    # 레이아웃 요약
    lines.append("### 사용 가능한 레이아웃")
    for layout in metadata["layouts"]:
        area = metadata["content_areas"].get(layout["name"], {})
        desc = area.get("description", "")
        lines.append(f"- **[{layout['index']}] {layout['name']}**: {desc}")
    lines.append("")

    # 콘텐츠 영역 상세
    lines.append("### 콘텐츠 영역 (content_areas)")
    lines.append("각 레이아웃에서 콘텐츠를 배치할 수 있는 영역:")
    content_area = metadata["content_areas"].get("내용", {})
    if content_area:
        header = content_area.get("header", {})
        body = content_area.get("body_area", {})
        lines.append(f"#### 내용 레이아웃 (주력)")
        lines.append(f"  - 문서명 (상단좌측): L={header.get('doc_title', {}).get('left')}, T={header.get('doc_title', {}).get('top')}, {header.get('doc_title', {}).get('font_size')}pt")
        lines.append(f"  - 목차경로 (breadcrumb): L={header.get('breadcrumb', {}).get('left')}, T={header.get('breadcrumb', {}).get('top')}, {header.get('breadcrumb', {}).get('font_size')}pt, Orange")
        lines.append(f"  - 메인제목: L={header.get('main_title', {}).get('left')}, T={header.get('main_title', {}).get('top')}, {header.get('main_title', {}).get('font_size')}pt, Bold")
        lines.append(f"  - 본문영역: L={body.get('left')}, T={body.get('top')}, W={body.get('width')}, H={body.get('height')}")
        lines.append(f"  - 섹션 타이틀: 12pt, Bold, Orange (본문 내 섹션 구분용)")
    lines.append("")

    # 스타일 가이드
    sg = metadata["style_guide"]
    lines.append("### 스타일 가이드")
    lines.append(f"- 기본 폰트: {sg['font_family']}")
    tc = sg["theme_colors"]
    lines.append(f"- 주요 색상: Blue={tc['accent1_blue']}, Orange={tc['accent2_orange']}, DarkGray={tc['dk2_dark_gray']}")
    ts = sg["table_style"]
    lines.append(f"- 테이블 헤더: fill={ts['header_fill']}, font=white/{ts['header_font_size']}pt/bold")
    lines.append(f"- 테이블 교대행: fill={ts['body_alt_fill']}")
    lines.append(f"- 테이블 다크네이비 헤더(대안): fill={ts['dark_navy_header']}")

    return "\n".join(lines)


def main():
    """CLI 진입점: 템플릿 인덱싱 실행"""
    import argparse

    parser = argparse.ArgumentParser(description="PPT 템플릿 인덱서")
    parser.add_argument(
        "--template",
        default=DEFAULT_TEMPLATE,
        help=f"템플릿 파일명 (기본: {DEFAULT_TEMPLATE})",
    )
    parser.add_argument(
        "--template-dir",
        default=str(TEMPLATE_DIR),
        help=f"템플릿 디렉토리 (기본: {TEMPLATE_DIR})",
    )
    args = parser.parse_args()

    template_dir = Path(args.template_dir)
    template_path = template_dir / args.template

    if not template_path.exists():
        print(f"[ERROR] 템플릿 파일을 찾을 수 없습니다: {template_path}")
        sys.exit(1)

    print(f"[INFO] 인덱싱 시작: {template_path}")

    metadata = index_template(template_path)

    # 결과 요약 출력
    layout_count = len(metadata["layouts"])
    print(f"[INFO] 슬라이드 크기: {metadata['slide_width_inches']}\" x {metadata['slide_height_inches']}\"")
    print(f"[INFO] 레이아웃 {layout_count}개 발견")

    for layout in metadata["layouts"]:
        shapes = layout.get("shapes", [])
        phs = layout.get("placeholders", [])
        print(f"  [{layout['index']}] {layout['name']}: Shape {len(shapes)}개, Placeholder {len(phs)}개")

    notes = metadata.get("analysis_notes", {})
    print(f"[INFO] Placeholder 기반: {notes.get('placeholder_based', 'N/A')}")
    print(f"[INFO] 기존 슬라이드: {notes.get('total_existing_slides', 0)}개")
    print(f"[INFO] 기존 테이블: {notes.get('table_count_in_existing', 0)}개")

    tc = metadata["style_guide"]["theme_colors"]
    print(f"[INFO] 테마 색상: accent1={tc['accent1_blue']}, accent2={tc['accent2_orange']}")

    # JSON 파일 저장
    output_path = template_dir / METADATA_FILENAME
    save_metadata(metadata, output_path)
    print(f"\n[OK] 메타데이터 저장 완료: {output_path}")

    # LLM용 요약 출력
    print("\n--- LLM 프롬프트용 요약 ---")
    print(format_metadata_for_llm(metadata))


if __name__ == "__main__":
    main()
